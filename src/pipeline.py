"""
Edu-Video AI Pipeline - Core Pipeline Functions
================================================
PDF processing, script generation, asset retrieval,
TTS/audio, video composition, PPTX generation.
"""
import os
import sys
import re
import subprocess
import json
import time
import gc
import tempfile
import shutil
import threading
import torch
import gradio as gr
import chromadb
import fitz
import requests
import numpy as np
from pathlib import Path
from transformers import AutoModel
from PIL import Image
from melo.api import TTS
from openvoice import se_extractor
from openvoice.api import ToneColorConverter
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import eval_engine

# Import all shared config
from config import (
    DEFAULT_AVATAR_ROOT, DEFAULT_AVATAR_ENV_PYTHON,
    DEFAULT_AVATAR_BACKEND, get_avatar_backend, build_avatar_command,
    DEFAULT_TTS_BACKEND, get_tts_backend,
    LLM_MODEL_PATH, DEFAULT_LLM_MODEL_PATH,
    LLM_ENV_PYTHON, LLM_WORKER_SCRIPT,
    _active_subprocesses, _cancel_autoprocess, _cancel_render,
    _layout_model, _layout_processor, _LAYOUT_MODEL_ID,
    _LAYOUT_VISUAL_LABELS, _LAYOUT_CAPTION_LABELS,
    _build_image_catalog, _load_layout_model, _detect_layout_regions,
    _cleanup_subprocess, _direct_slide_lookup,
    CHROMA_DB_PATH, METADATA_FOLDER, METADATA_PATH,
    create_placeholder_image, PLACEHOLDER_IMAGE,
    BROLL_ASSETS_DIR, PEXELS_API_KEY,
    CANVAS_PRESETS, PIP_SIZE_RATIO, PIP_PADDING,
    PROMPT_STANDARD, PROMPT_PROFESSOR, PROMPT_BROLL,
    PROMPT_PAPER, PROMPT_PAPER_BROLL, PROMPT_DECK,
    PROMPT_TECHNICAL_PAPER,
    IMAGE_CATALOG_TOP_K,
)

def generate_thumbnails_and_preview(pdf_file, progress=gr.Progress()):
    if not pdf_file:
        yield "❌ Please upload a PDF first.", gr.update(visible=False), gr.update(visible=False), [], [], []
        return
        
    yield "📄 Starting thumbnail generation...", gr.update(visible=False), gr.update(visible=False), [], [], []
    
    try:
        print(f"🔄 Processing PDF: {pdf_file.name}")
        doc = fitz.open(pdf_file.name)
        total_pages = len(doc)
        
        image_paths = []
        fullres_paths = []
        selected_indices = []
        
        thumb_dir = os.path.join(METADATA_FOLDER, "thumbnails")
        fullres_dir = os.path.join(METADATA_FOLDER, "fullres_pages")
        os.makedirs(thumb_dir, exist_ok=True)
        os.makedirs(fullres_dir, exist_ok=True)
        
        for i in range(total_pages):
            progress((i + 1) / total_pages, desc=f"Rendering page {i+1}/{total_pages}")
            page = doc[i]
            
            # Native resolution page render for gallery display (1x = PDF native DPI)
            pix = page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            thumb_path = os.path.join(thumb_dir, f"thumb_page_{i+1}.png")
            img.save(thumb_path, "PNG")
            
            # Full-resolution version for override usage (2x zoom, lossless PNG)
            fullres_pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            fullres_path = os.path.join(fullres_dir, f"fullres_page_{i+1}.png")
            fullres_pix.save(fullres_path)
            
            image_paths.append(thumb_path)
            fullres_paths.append(fullres_path)
            selected_indices.append(i) 
            
            yield f"⏳ Page {i+1}/{total_pages} done...", gr.update(), gr.update(), image_paths, selected_indices, fullres_paths
            
        print(f"✅ Generated {len(image_paths)} thumbnails + full-res pages successfully.")
        gallery_data = [(path, f"✅ Page {i+1}") for i, path in enumerate(image_paths)]
        
        yield (
            f"✅ {total_pages} thumbnails ready!",
            gr.update(value=gallery_data, visible=True),
            gr.update(visible=True),
            image_paths,
            selected_indices,
            fullres_paths
        )
    except Exception as e:
        print(f"❌ Thumbnail Error: {e}")
        yield f"❌ Error: {str(e)}", gr.update(visible=False), gr.update(visible=False), [], [], []

def bulk_select_all(image_paths_state):
    selected_indices = list(range(len(image_paths_state)))
    gallery_data = [(path, f"✅ Page {i+1}") for i, path in enumerate(image_paths_state)]
    return gallery_data, selected_indices

def bulk_deselect_all(image_paths_state):
    selected_indices = []
    gallery_data = [(path, f"❌ Page {i+1} (Skipped)") for i, path in enumerate(image_paths_state)]
    return gallery_data, selected_indices

def bulk_invert_selection(image_paths_state, selected_indices):
    new_selected = [i for i in range(len(image_paths_state)) if i not in selected_indices]
    gallery_data = [(path, f"✅ Page {i+1}" if i in new_selected else f"❌ Page {i+1} (Skipped)") for i, path in enumerate(image_paths_state)]
    return gallery_data, new_selected

def parse_page_ranges(range_str, max_pages):
    indices = set()
    if not range_str or not str(range_str).strip():
        return set()
    parts = [p.strip() for p in str(range_str).split(',')]
    for part in parts:
        if '-' in part:
            try:
                start, end = map(int, part.split('-'))
                start = max(1, start)
                end = min(max_pages, end)
                if start <= end:
                    indices.update(range(start - 1, end))
            except ValueError:
                continue
        else:
            try:
                val = int(part)
                if 1 <= val <= max_pages:
                    indices.add(val - 1)
            except ValueError:
                continue
    return indices

def select_range_fn(range_str, image_paths_state, selected_indices):
    target_indices = parse_page_ranges(range_str, len(image_paths_state))
    for idx in target_indices:
        if idx not in selected_indices:
            selected_indices.append(idx)
    gallery_data = [(path, f"✅ Page {i+1}" if i in selected_indices else f"❌ Page {i+1} (Skipped)") for i, path in enumerate(image_paths_state)]
    return gallery_data, selected_indices

def deselect_range_fn(range_str, image_paths_state, selected_indices):
    target_indices = parse_page_ranges(range_str, len(image_paths_state))
    selected_indices = [idx for idx in selected_indices if idx not in target_indices]
    gallery_data = [(path, f"✅ Page {i+1}" if i in selected_indices else f"❌ Page {i+1} (Skipped)") for i, path in enumerate(image_paths_state)]
    return gallery_data, selected_indices

def apply_preset(preset_name):
    if preset_name == "⚡ Fast Draft (Low VRAM)":
        return 4096, 512, 0.4
    elif preset_name == "📚 Full Lecture (High VRAM)":
        return 10240, 1024, 0.3
    else:
        return gr.update(), gr.update(), gr.update()

def update_prompt_text(choice):
    if choice == "Standard (Concise)":
        return PROMPT_STANDARD
    elif choice == "University Professor (Detailed)":
        return PROMPT_PROFESSOR
    elif choice == "Enhanced with B-Roll (PiP Layout)":
        return PROMPT_BROLL
    elif choice == "Academic Paper (Detailed)":
        return PROMPT_PAPER
    elif choice == "Academic Paper + B-Roll":
        return PROMPT_PAPER_BROLL
    elif choice == "Technical/STEM Paper":
        return PROMPT_TECHNICAL_PAPER
    else:
        return gr.update()
def _load_override_assets():
    """Load all available visual assets (extracted + custom) from metadata.json.
    Returns list of (image_path, label) tuples for the override gallery."""
    assets = []
    if os.path.exists(METADATA_PATH):
        try:
            with open(METADATA_PATH, "r", encoding="utf-8") as f:
                metadata = json.load(f)
            for i, item in enumerate(metadata):
                img_path = item.get("image_path", "")
                if img_path and os.path.exists(img_path):
                    caption = item.get("caption", "").strip()
                    page = item.get("page_number", "?")
                    if caption:
                        label = f"[{i}] P{page}: {caption[:50]}"
                    else:
                        label = f"[{i}] Page {page} asset"
                    assets.append((img_path, label))
        except Exception as e:
            print(f"⚠️ Failed to load override assets: {e}")
    return assets


def update_override_ui(script_json, pdf_images):
    """Populates the override dropdown with RETRIEVE_ASSET segments and the gallery
    with all extracted/custom visual assets from metadata.json (not just PDF pages)."""
    segments = extract_retrieve_segments(script_json)
    if not segments:
        return gr.update(choices=[], value=None), []
        
    choices = [f"Segment {s['segment_index'] + 1}: {s['query'][:40]}..." for s in segments]

    # Build gallery from extracted assets (metadata.json) — these are the actual
    # figures/tables/charts, not raw PDF pages
    gallery_items = _load_override_assets()

    # Fallback: if no extracted assets exist yet, show PDF page thumbnails
    if not gallery_items and pdf_images:
        gallery_items = [(img, f"Page {i+1}") for i, img in enumerate(pdf_images)]

    return gr.update(choices=choices, value=choices[0]), gallery_items

def handle_gallery_click(evt: gr.SelectData, selected_segment_str, current_state, pdf_images, pdf_fullres_images, current_gallery):
    """Triggered when a user clicks an asset thumbnail. Assigns it to the selected segment.
    Also live-updates the left-side asset_gallery to reflect the override."""
    if not selected_segment_str:
        return current_state, "⚠️ Please select a segment from the dropdown first!", current_gallery
        
    try:
        # Extract the segment number from the dropdown string (e.g., "Segment 2: ...")
        seg_num = int(selected_segment_str.split(":")[0].replace("Segment ", ""))
        seg_idx = seg_num - 1 # Convert to 0-based index
        
        # Try to resolve the clicked image from extracted assets first
        override_assets = _load_override_assets()
        if override_assets and evt.index < len(override_assets):
            selected_image_path = override_assets[evt.index][0]  # (path, label) tuple
            label = override_assets[evt.index][1]
        elif pdf_fullres_images and isinstance(pdf_fullres_images, list) and evt.index < len(pdf_fullres_images):
            selected_image_path = pdf_fullres_images[evt.index]
            label = f"Page {evt.index + 1}"
        elif pdf_images and evt.index < len(pdf_images):
            selected_image_path = pdf_images[evt.index]
            label = f"Page {evt.index + 1}"
        else:
            return current_state, "❌ Could not resolve the selected image.", current_gallery
        
        # Update the hidden state dictionary
        current_state[seg_idx] = selected_image_path
        
        # Live-update the left-side asset gallery: swap the matching segment's image
        updated_gallery = []
        seg_prefix = f"Segment {seg_num}:"
        replaced = False
        if current_gallery and isinstance(current_gallery, list):
            for item in current_gallery:
                # Gallery items are (path, label) tuples or similar
                if isinstance(item, (list, tuple)) and len(item) >= 2:
                    img, lbl = item[0], item[1]
                    if isinstance(lbl, str) and lbl.startswith(seg_prefix) and not replaced:
                        updated_gallery.append((selected_image_path, f"{seg_prefix} {label} [OVERRIDE]"))
                        replaced = True
                    else:
                        updated_gallery.append((img, lbl))
                else:
                    updated_gallery.append(item)
        
        if not replaced:
            updated_gallery = current_gallery  # No matching segment found, keep as-is
        
        return current_state, f"✅ Assigned {label} → Segment {seg_num}.", updated_gallery
    except Exception as e:
        return current_state, f"❌ Error: {e}", current_gallery

# --- Caption Detection Helper ---
_CAPTION_RE = re.compile(
    r'((?:Figure|Fig\.?|Table|Tab\.?|Chart|Diagram|Graph|Exhibit|Plate|Illustration|Equation|Eq\.?)\s*\.?\s*[\(]?\d+[a-z]?[\)]?[\.:)\-]?\s*[^\n]{0,300})',
    re.IGNORECASE
)

def _find_caption_near_rect(visual_rect, text_blocks, search_margin=120):
    """Search for figure/table captions in text blocks near a visual element.
    Returns (caption_text, caption_rect) or ("", None)."""
    search_area = visual_rect + fitz.Rect(-search_margin, -search_margin, search_margin, search_margin)
    best_caption, best_rect, best_dist = "", None, float("inf")

    for tb in text_blocks:
        if not search_area.intersects(tb["rect"]):
            continue
        match = _CAPTION_RE.search(tb["text"])
        if match:
            # Prefer captions closest to the visual (below first, then above)
            dist_below = abs(tb["rect"].y0 - visual_rect.y1)
            dist_above = abs(visual_rect.y0 - tb["rect"].y1)
            dist = min(dist_below, dist_above)
            if dist < best_dist:
                best_dist = dist
                best_caption = match.group(0).strip()
                best_rect = tb["rect"]

    return best_caption, best_rect

def _label_search_variants(label):
    """Generate lowercase search variants for 'Figure N' / 'Table N' / 'Equation N' labels.
    Handles abbreviations (Fig./Figure, Tab./Table, Eq./Equation) and case."""
    m = re.match(r'(Figure|Fig\.?|Table|Tab\.?|Chart|Diagram|Graph|Equation|Eq\.?)\s*\.?\s*[\(]?(\d+[a-z]?)[\)]?', label, re.IGNORECASE)
    if not m:
        return [label.lower()]
    word, num = m.group(1).lower().rstrip('.'), m.group(2)
    variants = set()
    if word in ('figure', 'fig'):
        for w in ['figure', 'fig.', 'fig']:
            variants.add(f"{w} {num}")
    elif word in ('table', 'tab'):
        for w in ['table', 'tab.', 'tab']:
            variants.add(f"{w} {num}")
    elif word in ('equation', 'eq'):
        for w in ['equation', 'eq.', 'eq']:
            variants.add(f"{w} {num}")
            variants.add(f"{w} ({num})")
        variants.add(f"equation ({num})")
    else:
        variants.add(f"{word} {num}")
    variants.add(label.lower())
    return list(variants)

def extract_universal_assets(pdf_path, output_folder, merge_distance=40):
    os.makedirs(output_folder, exist_ok=True)
    pdf_document = fitz.open(pdf_path)
    metadata_log = []
    
    first_page = pdf_document.load_page(0)
    is_presentation = first_page.rect.width > first_page.rect.height

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        raw_text = page.get_text("text").replace('\n', ' ').strip()
        semantic_context = f"Page/Slide {page_num + 1}: {raw_text}"
        
        if is_presentation:
            zoom = fitz.Matrix(2, 2)
            pixmap = page.get_pixmap(matrix=zoom)
            image_filename = f"slide_{page_num+1}.png"
            image_filepath = os.path.join(output_folder, image_filename)
            
            pixmap.save(image_filepath)
            if os.path.exists(image_filepath):
                metadata_log.append({
                    "image_path": image_filepath,
                    "page_number": page_num + 1,
                    "surrounding_text": semantic_context,
                    "caption": ""
                })
        else:
            # ── Build text-block index (always needed for caption search) ──
            page_dict = page.get_text("dict")
            text_blocks = []
            for b in page_dict.get("blocks", []):
                if b["type"] == 0:
                    block_text = ""
                    for line in b.get("lines", []):
                        for span in line.get("spans", []):
                            block_text += span.get("text", "")
                        block_text += " "
                    text_blocks.append({
                        "rect": fitz.Rect(b["bbox"]),
                        "text": block_text.strip()
                    })

            # ── Try AI-based layout detection first ──
            layout_detections = None
            try:
                # Render page to PIL for the layout model
                det_zoom = fitz.Matrix(2, 2)  # 2x for better detection
                det_pix = page.get_pixmap(matrix=det_zoom, alpha=False)
                page_pil = Image.frombytes("RGB", (det_pix.width, det_pix.height), det_pix.samples)
                layout_detections = _detect_layout_regions(page_pil, confidence=0.40)
            except Exception as det_err:
                print(f"⚠️ Layout detection failed on page {page_num+1}: {det_err}")

            if layout_detections is not None:
                # ── MODEL-BASED EXTRACTION ──
                pw, ph = page.rect.width, page.rect.height

                # Separate visual regions from caption regions
                visual_regions = []  # (fitz.Rect, label, confidence)
                caption_regions = [] # (fitz.Rect, label)
                for det in layout_detections:
                    label = det["label"]
                    # Convert from rendered-image pixel coords back to PDF coords
                    # det_zoom is 2x, so divide pixel coords by 2
                    x0_pdf = det["bbox"][0] / 2.0
                    y0_pdf = det["bbox"][1] / 2.0
                    x1_pdf = det["bbox"][2] / 2.0
                    y1_pdf = det["bbox"][3] / 2.0
                    rect = fitz.Rect(x0_pdf, y0_pdf, x1_pdf, y1_pdf)

                    # Clip to page bounds
                    rect = rect & page.rect

                    if label in _LAYOUT_VISUAL_LABELS:
                        # Filter out tiny noise detections
                        if rect.width > 20 and rect.height > 15:
                            visual_regions.append((rect, label, det["confidence"]))
                    elif label in _LAYOUT_CAPTION_LABELS:
                        caption_regions.append((rect, label))

                # ── Merge overlapping visual regions ──
                merged_rects = []
                merged_labels = []
                used = [False] * len(visual_regions)
                for i, (rect_i, label_i, conf_i) in enumerate(visual_regions):
                    if used[i]:
                        continue
                    current = fitz.Rect(rect_i)
                    current_label = label_i
                    used[i] = True
                    changed = True
                    while changed:
                        changed = False
                        for j, (rect_j, label_j, conf_j) in enumerate(visual_regions):
                            if used[j]:
                                continue
                            # Merge if overlapping or very close (within 15px)
                            expanded = current + (-15, -15, 15, 15)
                            if expanded.intersects(rect_j):
                                current = current | rect_j
                                used[j] = True
                                changed = True
                    merged_rects.append(current)
                    merged_labels.append(current_label)

                # ── Also grab PyMuPDF embedded image blocks as supplement ──
                pymupdf_image_rects = [fitz.Rect(b["bbox"]) for b in page_dict.get("blocks", []) if b["type"] == 1]
                for img_rect in pymupdf_image_rects:
                    if img_rect.width < 30 or img_rect.height < 30:
                        continue
                    # Check if already covered by a model detection
                    already_covered = any(
                        mr.intersects(img_rect) and (mr & img_rect).width > img_rect.width * 0.3
                        for mr in merged_rects
                    )
                    if not already_covered:
                        merged_rects.append(img_rect)
                        merged_labels.append("Picture")

                # ── Associate captions and render ──
                for cluster_index, rect in enumerate(merged_rects):
                    label = merged_labels[cluster_index]
                    is_equation = (label == "Formula")
                    min_size = 15 if is_equation else 50

                    if rect.width > min_size and rect.height > min_size:
                        # Search for caption near this region
                        caption, caption_rect = _find_caption_near_rect(rect, text_blocks, search_margin=80)

                        # For captions detected by model, also check caption_regions
                        if not caption:
                            for cap_rect, _ in caption_regions:
                                # Caption must be near and horizontally overlapping
                                dist_below = cap_rect.y0 - rect.y1
                                dist_above = rect.y0 - cap_rect.y1
                                h_overlap = min(rect.x1, cap_rect.x1) - max(rect.x0, cap_rect.x0)
                                if h_overlap > 20 and (-10 < dist_below < 60 or -10 < dist_above < 60):
                                    # Read the caption text from text_blocks
                                    for tb in text_blocks:
                                        if cap_rect.intersects(tb["rect"]):
                                            m = _CAPTION_RE.search(tb["text"])
                                            if m:
                                                caption = m.group(0).strip()
                                                caption_rect = cap_rect
                                                break
                                    if caption:
                                        break

                        # Build clip rect (include caption in crop)
                        clip_rect = fitz.Rect(rect)
                        if caption_rect:
                            clip_rect = clip_rect | caption_rect

                        zoom = fitz.Matrix(2, 2)
                        pixmap = page.get_pixmap(matrix=zoom, clip=clip_rect)

                        image_filename = f"page_{page_num+1}_cluster_{cluster_index}.png"
                        image_filepath = os.path.join(output_folder, image_filename)

                        pixmap.save(image_filepath)
                        if os.path.exists(image_filepath):
                            metadata_log.append({
                                "image_path": image_filepath,
                                "page_number": page_num + 1,
                                "surrounding_text": semantic_context,
                                "caption": caption
                            })

            else:
                # ── FALLBACK: Geometric heuristic extraction ──
                # (used when layout model is not available)

                # Collect image blocks
                visual_rects = [fitz.Rect(b["bbox"]) for b in page_dict.get("blocks", []) if b["type"] == 1]

                # Detect text-based tables (PyMuPDF >= 1.23)
                try:
                    tables = page.find_tables()
                    for tab in tables.tables:
                        table_rect = fitz.Rect(tab.bbox)
                        if table_rect.width > 50 and table_rect.height > 50:
                            visual_rects.append(table_rect)
                except (AttributeError, Exception):
                    pass

                # Merge nearby visual elements
                h_merge = 10
                v_merge = merge_distance
                merged_rects = []
                while visual_rects:
                    current_rect = visual_rects.pop(0)
                    merged = True
                    while merged:
                        merged = False
                        for i, other_rect in enumerate(visual_rects):
                            search_area = current_rect + (-h_merge, -v_merge, h_merge, v_merge)
                            if search_area.intersects(other_rect):
                                current_rect = current_rect | other_rect
                                visual_rects.pop(i)
                                merged = True
                                break
                    merged_rects.append(current_rect)

                # Caption-based table fallback — scan BELOW and ABOVE caption
                for tb in text_blocks:
                    if not re.match(r'\s*(Table|Tab\.?)\s*\.?\s*\d+', tb["text"], re.IGNORECASE):
                        continue
                    cap_rect = tb["rect"]

                    # Check if already covered
                    search_zone = fitz.Rect(
                        cap_rect.x0 - 30, cap_rect.y0 - 300,
                        cap_rect.x1 + 30, cap_rect.y1 + 300
                    )
                    already_covered = any(
                        search_zone.intersects(mr) and mr.height > 40
                        for mr in merged_rects
                    )
                    if already_covered:
                        continue

                    # --- Try BELOW caption first (IEEE/ACM style: caption above table) ---
                    blocks_below = []
                    for b in text_blocks:
                        if b is tb:
                            continue
                        if b["rect"].y0 <= cap_rect.y1 + 5:
                            continue  # Must be below the caption
                        horiz_overlap = min(cap_rect.x1, b["rect"].x1) - max(cap_rect.x0, b["rect"].x0)
                        if horiz_overlap > 10:
                            blocks_below.append(b)
                    blocks_below = sorted(blocks_below, key=lambda b: b["rect"].y0)  # closest first

                    table_bottom = cap_rect.y1 + 20  # default
                    prev_y1 = cap_rect.y1
                    found_below = False
                    for below_tb in blocks_below:
                        gap = below_tb["rect"].y0 - prev_y1
                        if gap > 60:
                            table_bottom = below_tb["rect"].y0 - 2
                            found_below = True
                            break
                        table_bottom = below_tb["rect"].y1 + 5
                        prev_y1 = below_tb["rect"].y1
                    
                    if found_below and table_bottom - cap_rect.y1 > 30:
                        table_region = fitz.Rect(
                            cap_rect.x0 - 10, cap_rect.y0 - 5,
                            cap_rect.x1 + 10, table_bottom
                        )
                        if table_region.width > 50 and table_region.height > 30:
                            merged_rects.append(table_region)
                            continue

                    # --- Try ABOVE caption (classic style: caption below table) ---
                    blocks_above = []
                    for b in text_blocks:
                        if b is tb:
                            continue
                        if b["rect"].y1 >= cap_rect.y0 - 5:
                            continue
                        horiz_overlap = min(cap_rect.x1, b["rect"].x1) - max(cap_rect.x0, b["rect"].x0)
                        if horiz_overlap > 10:
                            blocks_above.append(b)
                    blocks_above = sorted(blocks_above, key=lambda b: b["rect"].y1, reverse=True)

                    table_top = cap_rect.y0 - 20
                    prev_y0 = cap_rect.y0
                    for above_tb in blocks_above:
                        gap = prev_y0 - above_tb["rect"].y1
                        if gap > 60:
                            table_top = above_tb["rect"].y1 + 2
                            break
                        table_top = above_tb["rect"].y0 - 5
                        prev_y0 = above_tb["rect"].y0
                    else:
                        table_top = max(page.rect.y0 + 10, cap_rect.y0 - 350)

                    table_region = fitz.Rect(
                        cap_rect.x0 - 10, table_top,
                        cap_rect.x1 + 10, cap_rect.y1 + 10
                    )
                    if table_region.width > 50 and table_region.height > 30:
                        merged_rects.append(table_region)

                # Equation detection
                _eq_num_end_re = re.compile(r'\(\s*(\d+[a-z]?)\s*\)\s*$')
                _eq_num_standalone_re = re.compile(r'^\s*\((\d+[a-z]?)\)\s*$')
                equation_caption_map = {}
                used_eq_blocks = set()

                page_right = page.rect.x1
                page_left = page.rect.x0
                page_cx = (page_left + page_right) / 2
                page_w = page.rect.width
                right_margin_threshold = page_right - page_w * 0.20

                eq_anchors = []
                for idx_tb, tb_item in enumerate(text_blocks):
                    text_s = tb_item["text"].strip()
                    block_rect = tb_item["rect"]
                    if any(mr.intersects(block_rect) for mr in merged_rects):
                        continue
                    m_standalone = _eq_num_standalone_re.match(text_s)
                    if m_standalone and block_rect.x0 > page_cx:
                        eq_anchors.append((m_standalone.group(1), block_rect, idx_tb))
                        continue
                    if len(text_s) > 400:
                        continue
                    m_end = _eq_num_end_re.search(text_s)
                    if m_end and block_rect.x1 >= right_margin_threshold:
                        eq_anchors.append((m_end.group(1), block_rect, idx_tb))

                for eq_num, anchor_rect, anchor_idx in eq_anchors:
                    used_eq_blocks.add(anchor_idx)
                    eq_region = fitz.Rect(anchor_rect)
                    for idx_tb2, other_tb in enumerate(text_blocks):
                        if idx_tb2 == anchor_idx or idx_tb2 in used_eq_blocks:
                            continue
                        other_rect = other_tb["rect"]
                        other_text = other_tb["text"].strip()
                        v_overlap = min(eq_region.y1, other_rect.y1) - max(eq_region.y0, other_rect.y0)
                        gap_above = eq_region.y0 - other_rect.y1
                        gap_below = other_rect.y0 - eq_region.y1
                        is_near = v_overlap > -5 or (0 < gap_above < 30) or (0 < gap_below < 30)
                        if not is_near:
                            continue
                        if len(other_text) > 250:
                            continue
                        if other_rect.x0 < page_left + page_w * 0.08 and other_rect.x1 > page_right - page_w * 0.08:
                            if len(other_text) > 80:
                                continue
                        eq_region = eq_region | other_rect
                        used_eq_blocks.add(idx_tb2)
                    eq_clip = fitz.Rect(
                        page_left + 25, eq_region.y0 - 8,
                        page_right - 25, eq_region.y1 + 8
                    )
                    if eq_clip.width > 40 and eq_clip.height > 12:
                        idx = len(merged_rects)
                        merged_rects.append(eq_clip)
                        equation_caption_map[idx] = f"Equation ({eq_num})"

                # Render each cluster
                for cluster_index, rect in enumerate(merged_rects):
                    is_equation = cluster_index in equation_caption_map
                    min_size = 15 if is_equation else 50
                    if rect.width > min_size and rect.height > min_size:
                        if is_equation:
                            caption = equation_caption_map[cluster_index]
                            clip_rect = rect
                        else:
                            caption, caption_rect = _find_caption_near_rect(rect, text_blocks)
                            clip_rect = rect
                            if caption_rect:
                                clip_rect = clip_rect | caption_rect
                        zoom = fitz.Matrix(2, 2)
                        pixmap = page.get_pixmap(matrix=zoom, clip=clip_rect)
                        image_filename = f"page_{page_num+1}_cluster_{cluster_index}.png"
                        image_filepath = os.path.join(output_folder, image_filename)
                        pixmap.save(image_filepath)
                        if os.path.exists(image_filepath):
                            metadata_log.append({
                                "image_path": image_filepath,
                                "page_number": page_num + 1,
                                "surrounding_text": semantic_context,
                                "caption": caption
                            })

    metadata_path = os.path.join(output_folder, "metadata.json")
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(metadata_log, f, indent=4)
        
    return metadata_path

def _proportional_sample_pages(raw_text: str, budget: int) -> str:
    """
    Distribute the char budget proportionally across all page sections so that
    every page is represented in the LLM context instead of only the first N chars.
    Falls back to a plain slice when no page markers are found.
    """
    parts = re.split(r'(\n--- PAGE \d+ ---\n)', raw_text)

    pages: list[tuple[str, str]] = []  # (header, content)
    i = 0
    while i < len(parts):
        if re.match(r'\n--- PAGE \d+ ---\n', parts[i]):
            header = parts[i]
            content = parts[i + 1] if (i + 1) < len(parts) else ""
            pages.append((header, content))
            i += 2
        else:
            if parts[i].strip():
                pages.append(("", parts[i]))
            i += 1

    if not pages:
        return raw_text[:budget]

    total_content = sum(len(c) for _, c in pages)
    if total_content == 0:
        return raw_text[:budget]

    result: list[str] = []
    for header, content in pages:
        # Each page gets a slice proportional to its share of the document,
        # but never more than the page actually contains and at least 100 chars.
        alloc = min(len(content), max(100, int(len(content) / total_content * budget)))
        result.append(header + content[:alloc])

    return "".join(result)[:budget]  # safety cap


# ─── System prompt used during the MAP phase of map-reduce ──────────────────────
_MAP_EXTRACT_PROMPT = """You are a precise document analyst. Extract the essential information from this document section.

Capture ALL of the following in your response:
- Main claims, arguments, or findings
- Methods, algorithms, or procedures described
- Quantitative results, metrics, or numbers mentioned
- Key concepts, definitions, or terminology introduced
- Figures, tables, equations, or diagrams referenced (with brief descriptions of what they show)

Output ONLY valid JSON (no extra text) in this exact format:
{
  "section_summary": "2-3 sentence summary of this section",
  "key_points": ["point 1", "point 2"],
  "figures_referenced": ["Figure 1: description of what it shows"],
  "quantitative_results": ["result/metric 1"]
}"""


def _run_llm_blocking(system_prompt: str, user_content: str, n_ctx: int, n_batch: int,
                      temperature: float, active_model: str) -> "str | None":
    """
    Run the LLM worker subprocess in non-streaming (blocking) mode.
    Returns the full response string, or None on failure.
    Used by the sliding-window and map-reduce strategies.
    """
    llm_config = {
        "model_path": active_model,
        "system_prompt": system_prompt,
        "user_content": user_content,
        "n_ctx": int(n_ctx),
        "n_batch": int(n_batch),
        "temperature": float(temperature),
        "max_tokens": -1,
        "top_p": 0.95,
        "stream": False,
        "json_mode": True,
        "stop": ["</s>", "<end_of_turn>", "<eos>"]
    }
    config_fd, config_path = tempfile.mkstemp(suffix=".json", prefix="llm_cfg_")
    try:
        with os.fdopen(config_fd, "w", encoding="utf-8") as cf:
            json.dump(llm_config, cf)
        proc = subprocess.run(
            [LLM_ENV_PYTHON, LLM_WORKER_SCRIPT, config_path],
            capture_output=True, text=True, encoding="utf-8", timeout=1800
        )
        for line in proc.stdout.splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                data = json.loads(line)
                if "content" in data:
                    return data["content"]
                elif "error" in data:
                    print(f"\u26a0\ufe0f LLM blocking call error: {data['error']}")
                    return None
            except json.JSONDecodeError:
                continue
    except subprocess.TimeoutExpired:
        print("\u26a0\ufe0f LLM blocking call timed out after 30 minutes")
    except Exception as e:
        print(f"\u26a0\ufe0f LLM blocking call failed: {e}")
    finally:
        try:
            os.unlink(config_path)
        except OSError:
            pass
    return None


def _split_into_page_chunks(raw_text: str, budget: int, overlap_pages: int = 1) -> "list[str]":
    """
    Split raw_text into chunks that each fit within `budget` chars.
    Splits along --- PAGE N --- boundaries (never mid-page).
    `overlap_pages`: number of trailing pages from the previous chunk
    prepended to the next chunk for continuity.
    Falls back to plain char-slicing when no page markers exist.
    """
    parts = re.split(r'(\n--- PAGE \d+ ---\n)', raw_text)
    pages: list = []
    i = 0
    while i < len(parts):
        if re.match(r'\n--- PAGE \d+ ---\n', parts[i]):
            header = parts[i]
            content = parts[i + 1] if (i + 1) < len(parts) else ""
            pages.append(header + content)
            i += 2
        else:
            if parts[i].strip():
                pages.append(parts[i])
            i += 1

    if not pages:
        return [raw_text[off:off + budget] for off in range(0, len(raw_text), budget)]

    chunks: list = []
    current_pages: list = []
    current_len = 0
    for page in pages:
        if current_len + len(page) > budget and current_pages:
            chunks.append("".join(current_pages))
            # Carry last N pages into the next chunk for context continuity
            current_pages = current_pages[-overlap_pages:] if overlap_pages else []
            current_len = sum(len(p) for p in current_pages)
        current_pages.append(page)
        current_len += len(page)
    if current_pages:
        chunks.append("".join(current_pages))

    return chunks


def _rank_catalog_by_relevance(metadata_path: str, query_text: str, top_k: int):
    """Rank catalog entries by JinaCLIP cosine similarity against `query_text`.

    Returns a list of original metadata indices, ordered most-relevant first,
    capped at `top_k`. If the catalog is already <= top_k, or if anything fails
    (model load error, empty query, etc.), returns None to signal the caller
    should emit the full catalog unchanged.

    Why this exists: long papers can produce 60–100+ extracted assets. Listing
    every single one (with caption + page) in the LLM prompt can blow past the
    10,240-token context window and silently push earlier prompt content out.
    Pre-ranking trims the catalog to the most semantically relevant entries
    while keeping the original `[Image N]` IDs intact, so downstream retrieval
    still works for any image (the LLM just won't see the long tail unless it
    happens to invent a matching query).
    """
    try:
        if not query_text or not query_text.strip():
            return None
        if not os.path.exists(metadata_path):
            return None
        with open(metadata_path, "r", encoding="utf-8") as f:
            metadata = json.load(f)
        if not metadata or len(metadata) <= top_k:
            return None  # Catalog already fits; no ranking needed

        # Build the document-side text for each entry (caption preferred, surrounding fallback)
        entry_texts = []
        for item in metadata:
            cap = (item.get("caption") or "").strip()
            if not cap:
                cap = (item.get("surrounding_text") or "").strip()[:400]
            entry_texts.append(cap if cap else "uncaptioned visual")

        device = "cuda" if torch.cuda.is_available() else "cpu"
        model = AutoModel.from_pretrained(
            'jinaai/jina-clip-v1', trust_remote_code=True,
            low_cpu_mem_usage=False, _fast_init=False,
        ).to(device)

        # Cap query to a sensible length so the encoder stays fast
        q_text = query_text.strip()[:1500]
        q_emb = model.encode_text([q_text])  # shape (1, D)
        e_emb = model.encode_text(entry_texts)  # shape (N, D)

        # Cosine similarity (encode_text outputs are already normalised by JinaCLIP)
        import numpy as _np
        q = _np.asarray(q_emb, dtype=_np.float32)
        e = _np.asarray(e_emb, dtype=_np.float32)
        # Defensive renormalisation
        q = q / (_np.linalg.norm(q, axis=1, keepdims=True) + 1e-8)
        e = e / (_np.linalg.norm(e, axis=1, keepdims=True) + 1e-8)
        scores = (e @ q.T).reshape(-1)  # (N,)

        # Take top-k indices, descending by score
        order = _np.argsort(-scores)[:top_k]
        return [int(i) for i in order]
    except Exception as e:
        print(f"⚠️ Catalog pre-ranking failed ({e}); falling back to full catalog.")
        return None


def step1_generate_script(pdf_file, selected_indices, ui_n_ctx, ui_n_batch, ui_temp, custom_system_prompt, focus_areas="", llm_model_path="", long_doc_mode="Auto (Sliding Window / Map-Reduce)"):
    if not pdf_file or not selected_indices:
        yield "❌ Error: Please upload a PDF and select at least one page.", ""
        return

    yield "📄 Extracting text from selected pages...", ""
    
    # Resolve LLM model path (allow UI override)
    active_model = llm_model_path.strip() if llm_model_path and llm_model_path.strip() else LLM_MODEL_PATH
    if not os.path.exists(active_model):
        yield f"❌ LLM model not found at: {active_model}", ""
        return
    
    try:
        doc = fitz.open(pdf_file.name)
        raw_text = ""
        for i in range(len(doc)):
            if i in selected_indices:
                raw_text += f"\n--- PAGE {i+1} ---\n" + doc[i].get_text() + "\n"
        
        safe_char_limit = int((ui_n_ctx * 3.5) * 0.40)

        # ── Strategy selection based on how severely the doc exceeds budget ──
        ratio = len(raw_text) / safe_char_limit if safe_char_limit > 0 else float("inf")
        n_chunks_est = max(1, (len(raw_text) + safe_char_limit - 1) // safe_char_limit)

        force_fast = "fast" in long_doc_mode.lower() or "truncat" in long_doc_mode.lower()

        if ratio <= 1.30 or force_fast:
            # Fits, or user explicitly chose fast/truncate mode
            strategy = "single"
            chunk = _proportional_sample_pages(raw_text, safe_char_limit) if (ratio > 1.30 and force_fast) else raw_text[:safe_char_limit]
            if ratio > 1.30 and force_fast:
                pct_lost = round((1 - safe_char_limit / len(raw_text)) * 100, 1)
                trunc_warning = (
                    f"⚠️ Fast mode: document is {ratio:.1f}x the context budget "
                    f"({len(raw_text):,} chars vs {safe_char_limit:,} budget, ~{pct_lost}% lost). "
                    f"Using proportional page sampling (single pass)."
                )
            else:
                trunc_warning = None
        elif ratio <= 5.0:
            # Moderately large (1.3x–5x): sliding window
            strategy = "sliding_window"
            chunk = None
            trunc_warning = (
                f"⚠️ Document is {ratio:.1f}x the context budget "
                f"({len(raw_text):,} chars vs {safe_char_limit:,} budget). "
                f"Using sliding-window strategy (~{n_chunks_est} passes with 1-page overlap) "
                f"— all content will be processed and segments merged."
            )
        else:
            # Very large (>5x): map-reduce
            strategy = "map_reduce"
            chunk = None
            trunc_warning = (
                f"⚠️ Document is {ratio:.1f}x the context budget "
                f"({len(raw_text):,} chars vs {safe_char_limit:,} budget). "
                f"Using map-reduce strategy: extracting facts from {n_chunks_est} sections, "
                f"then synthesising the final script in one reduce pass."
            )
        # ─────────────────────────────────────────────────────────────────────
    except Exception as e:
        yield f"❌ Error reading PDF: {e}", ""
        return

    if trunc_warning:
        yield trunc_warning, ""

    yield f"🧠 Booting up Semantic Engine...", ""

    # ── Build image catalog from PDF-native data (captions + surrounding text) ──
    image_catalog = ""
    try:
        if not os.path.exists(METADATA_PATH) and pdf_file:
            yield "🖼️ Extracting visual assets from PDF for image catalog...", ""
            extract_universal_assets(pdf_file.name, METADATA_FOLDER)
        if os.path.exists(METADATA_PATH):
            # Pre-rank with JinaCLIP when the catalog is large enough that listing
            # every entry would risk overflowing the LLM context window. The query
            # is the document overview (first ~1500 chars of raw_text). Original
            # [Image N] IDs are preserved so retrieval still works for any image.
            doc_overview = (raw_text or "").strip()[:1500]
            ranked = _rank_catalog_by_relevance(METADATA_PATH, doc_overview, IMAGE_CATALOG_TOP_K)
            if ranked is not None:
                try:
                    with open(METADATA_PATH, "r", encoding="utf-8") as _f:
                        _total = len(json.load(_f))
                except Exception:
                    _total = len(ranked)
                yield f"🎯 Pre-ranking image catalog: kept top {len(ranked)} of {_total} extracted assets to fit the context window...", ""
            image_catalog = _build_image_catalog(METADATA_PATH, selected_indices=ranked, ranked_note=True)
    except Exception as e:
        print(f"⚠️ Image catalog generation failed ({e}), proceeding without catalog")

    # ── Inline helper: attach image catalog + focus instructions ──────────────
    def _attach_catalog_and_focus(uc, text_src):
        """Append image catalog (with re-trim) and focus-area instructions."""
        if image_catalog:
            uc += f"\n\n{image_catalog}"
            total_limit = int((ui_n_ctx * 3.5) * 0.40)
            if len(uc) > total_limit:
                cat_len = len(image_catalog) + 2
                remaining = total_limit - cat_len
                if remaining > 500:
                    uc = f"Text to convert:\n{_proportional_sample_pages(text_src, remaining)}\n\n{image_catalog}"
        if focus_areas and str(focus_areas).strip():
            uc += (
                f"\n\nSPECIAL INSTRUCTION FROM USER:\n"
                f"The audience specifically wants extra focus and elaboration on these areas:\n"
                f"{focus_areas.strip()}\n\n"
                f"When generating the JSON script:\n"
                f"- Create richer, longer spoken_text segments (3–6 sentences) for these topics.\n"
                f"- Add more educational examples, analogies, or why-it-matters explanations.\n"
                f"- Feel free to create additional segments if needed to cover these areas properly.\n"
                f"- Keep other segments concise as usual."
            )
        return uc

    if strategy == "sliding_window":
        # ── SLIDING WINDOW ────────────────────────────────────────────────────
        # Splits the document into page-boundary chunks, runs the full script-
        # generation prompt on each chunk, then merges all segments arrays.
        # Each chunk gets a position-aware directive appended to the system
        # prompt so the LLM doesn't write a fresh intro/outro per pass (which
        # otherwise causes duplicated "Welcome..." and "In conclusion..."
        # segments after merging).
        # overlap_pages=0: the position directive already gives the LLM enough
        # cross-chunk context. Re-narrating the boundary page caused duplicate
        # figure walkthroughs in earlier runs without enough context benefit
        # to justify it (the merge step does no overlap resolution).
        chunks = _split_into_page_chunks(raw_text, safe_char_limit, overlap_pages=0)
        all_segments = []
        n_chunks = len(chunks)
        # Forbidden opener phrases for middle/final chunks. The LLM tends to
        # paraphrase around a simple "no Welcome" rule (e.g. "We are diving
        # into…", "Our objective…"), so we enumerate the common workarounds.
        forbidden_openers = (
            '"Welcome", "Today we are examining", "Today we\'re examining", '
            '"Today, we", "We are diving", "We\'re diving", "We are now moving", '
            '"We will explore", "We\'ll explore", "Let\'s dive", "Let\'s explore", '
            '"Let\'s look at", "Our objective", "In this work", "In this paper", '
            '"In this section", "Next, we will", "Now, let\'s", "First, let\'s", '
            '"This presentation", "This lecture", "This deep dive"'
        )
        for idx, text_chunk in enumerate(chunks):
            yield f"🔄 Sliding window pass {idx + 1}/{n_chunks}: generating script for page group...", ""
            uc = _attach_catalog_and_focus(f"Text to convert:\n{text_chunk}", text_chunk)

            # Position-aware directive: tell the LLM where this chunk sits in
            # the overall narration so it produces a single coherent script
            # across passes rather than N standalone scripts.
            if n_chunks == 1:
                position_directive = ""
            elif idx == 0:
                position_directive = (
                    "\n\n=== SLIDING-WINDOW POSITION: OPENING SECTION (1 of "
                    f"{n_chunks}) ===\n"
                    "This is the FIRST chunk of a multi-pass narration. "
                    "Begin with a welcome/introduction segment as usual. "
                    "Do NOT write a conclusion or wrap-up — more chunks will follow. "
                    "End the segments array on a natural mid-narrative note."
                )
            elif idx == n_chunks - 1:
                position_directive = (
                    f"\n\n=== SLIDING-WINDOW POSITION: FINAL SECTION ({idx + 1} of "
                    f"{n_chunks}) ===\n"
                    "This is the LAST chunk of a multi-pass narration. "
                    "The audience has ALREADY heard a welcome/introduction in an earlier pass. "
                    "The FIRST segment must NOT begin with any of these FORBIDDEN OPENERS:\n"
                    f"{forbidden_openers}.\n"
                    "Instead, start the first segment with a continuation phrase such as "
                    '"Building on the previous section,", "Another important aspect is",  '
                    '"Turning to", "Beyond this,", or directly with the figure/concept name '
                    '(e.g. "Figure 8 shows…", "The fine-tuning procedure…"). '
                    "End with a brief conclusion segment that wraps up the overall paper."
                )
            else:
                position_directive = (
                    f"\n\n=== SLIDING-WINDOW POSITION: MIDDLE SECTION ({idx + 1} of "
                    f"{n_chunks}) ===\n"
                    "This is a MIDDLE chunk of a multi-pass narration. "
                    "The audience has ALREADY heard a welcome/introduction in an earlier pass, "
                    "and more chunks will follow this one. "
                    "The FIRST segment must NOT begin with any of these FORBIDDEN OPENERS:\n"
                    f"{forbidden_openers}.\n"
                    "Instead, start the first segment with a continuation phrase such as "
                    '"Building on the previous section,", "Another important aspect is", '
                    '"Turning to", "Beyond this,", or directly with the figure/concept name '
                    '(e.g. "Figure 5 shows…", "The text embeddings are extracted…"). '
                    "Do NOT write a conclusion or 'In conclusion...' segment — more chunks follow. "
                    "Treat this as a mid-stream continuation, not a standalone script."
                )

            chunk_system_prompt = custom_system_prompt + position_directive
            raw_result = _run_llm_blocking(
                chunk_system_prompt, uc, ui_n_ctx, ui_n_batch, ui_temp, active_model
            )
            if raw_result:
                raw_result = raw_result.strip()
                if not raw_result.endswith("}"):
                    raw_result += "\n}"
                try:
                    parsed = json.loads(raw_result)
                    segs = parsed.get("segments", [])
                    all_segments.extend(segs)
                    yield f"✅ Pass {idx + 1}/{len(chunks)} complete — {len(segs)} new segments", ""
                except json.JSONDecodeError:
                    yield f"⚠️ Pass {idx + 1}/{len(chunks)} JSON parse failed, skipping", ""
            else:
                yield f"⚠️ Pass {idx + 1}/{len(chunks)} returned empty, skipping", ""
        if not all_segments:
            yield "❌ Sliding window: no segments generated across all passes", ""
            return

        # ── Post-merge dedup: drop later RETRIEVE_ASSET segments whose query
        # already appeared earlier (case-insensitive, whitespace-normalised).
        # Catches duplicate figure/table walkthroughs caused by the same asset
        # being narrated by multiple chunks (e.g. when a figure caption spans
        # multiple pages, or when middle chunks re-introduce a figure shown
        # earlier). Keeps GENERATE_AVATAR segments untouched.
        seen_queries: set = set()
        deduped_segments: list = []
        dropped = 0
        for seg in all_segments:
            mode = (seg.get("visual_mode") or "").upper()
            if mode == "RETRIEVE_ASSET":
                q = " ".join((seg.get("query") or "").lower().split())
                # Use a prefix key so paraphrased queries pointing at the same
                # figure (e.g. truncated vs full caption) collapse together.
                key = q[:120]
                if key and key in seen_queries:
                    dropped += 1
                    continue
                if key:
                    seen_queries.add(key)
            deduped_segments.append(seg)
        if dropped:
            yield f"🧹 Deduped {dropped} repeat asset segment(s) across passes", ""

        merged = json.dumps({"segments": deduped_segments}, ensure_ascii=False, indent=2)
        yield (
            f"✅ Script generation complete! "
            f"({len(deduped_segments)} segments from {n_chunks} passes, "
            f"{dropped} duplicate assets removed)"
        ), merged
        return

    elif strategy == "map_reduce":
        # ── MAP-REDUCE ────────────────────────────────────────────────────────
        # MAP: ask the LLM to extract structured facts from each page-group.
        # REDUCE: feed all extracted facts to the LLM to generate the script.
        chunks = _split_into_page_chunks(raw_text, safe_char_limit)
        map_results = []
        for idx, text_chunk in enumerate(chunks):
            yield f"🗺️ Map phase {idx + 1}/{len(chunks)}: extracting key facts...", ""
            raw_map = _run_llm_blocking(
                _MAP_EXTRACT_PROMPT,
                f"Document section:\n{text_chunk}",
                ui_n_ctx, ui_n_batch, ui_temp, active_model
            )
            if raw_map:
                try:
                    map_results.append(json.loads(raw_map))
                    yield f"✅ Map {idx + 1}/{len(chunks)} complete", ""
                except json.JSONDecodeError:
                    map_results.append({"raw_facts": raw_map[:800]})
                    yield f"⚠️ Map {idx + 1}/{len(chunks)} JSON parse failed, using raw text", ""
            else:
                yield f"⚠️ Map {idx + 1}/{len(chunks)} returned empty, skipping", ""
        if not map_results:
            yield "❌ Map-reduce: map phase produced no results", ""
            return
        # Reduce: build user_content from extracted facts, then stream final script
        facts_json = json.dumps({"document_sections": map_results}, ensure_ascii=False)
        cat_reserve = len(image_catalog) + 2 if image_catalog else 0
        text_budget = max(500, safe_char_limit - cat_reserve)
        if len(facts_json) > text_budget:
            facts_json = facts_json[:text_budget]
        user_content = (
            f"Document facts extracted from {len(chunks)} sections "
            f"(synthesise a complete script covering all sections):\n{facts_json}"
        )
        user_content = _attach_catalog_and_focus(user_content, facts_json)
        yield f"🔀 Reduce phase: synthesising final script from {len(chunks)} section summaries...", ""

    else:
        # ── SINGLE PASS ───────────────────────────────────────────────────────
        user_content = _attach_catalog_and_focus(f"Text to convert:\n{chunk}", chunk)

    # ── Streaming LLM call (single-pass and map-reduce reduce phase) ──────────
    try:
        yield f"⏳ Launching LLM in isolated environment (n_ctx={ui_n_ctx}, batch={ui_n_batch})...", ""

        # Write config for the LLM worker subprocess
        llm_config = {
            "model_path": active_model,
            "system_prompt": custom_system_prompt,
            "user_content": user_content,
            "n_ctx": int(ui_n_ctx),
            "n_batch": int(ui_n_batch),
            "temperature": float(ui_temp),
            "max_tokens": -1,
            "top_p": 0.95,
            "stream": True,
            "json_mode": True,
            "stop": ["</s>", "<end_of_turn>", "<eos>"]
        }

        config_fd, config_path = tempfile.mkstemp(suffix=".json", prefix="llm_cfg_")
        try:
            with os.fdopen(config_fd, "w", encoding="utf-8") as cf:
                json.dump(llm_config, cf)

            if not os.path.exists(LLM_ENV_PYTHON):
                yield f"❌ LLM environment not found at: {LLM_ENV_PYTHON}\nRun: pip install llama-cpp-python in the llm_env venv.", ""
                return

            yield "✍️ Generating JSON script...", ""

            process = subprocess.Popen(
                [LLM_ENV_PYTHON, LLM_WORKER_SCRIPT, config_path],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                bufsize=1,
                encoding="utf-8"
            )
            _active_subprocesses["script_llm"] = process

            script_text = ""
            for line in process.stdout:
                line = line.strip()
                if not line:
                    continue
                try:
                    data = json.loads(line)
                    if "t" in data:
                        script_text += data["t"]
                        yield "✍️ Generating JSON script (Streaming...)", script_text
                    elif "error" in data:
                        yield f"❌ LLM Error: {data['error']}", ""
                        process.kill()
                        _active_subprocesses.pop("script_llm", None)
                        return
                    elif "done" in data:
                        break
                except json.JSONDecodeError:
                    continue

            _active_subprocesses.pop("script_llm", None)
            process.wait()
            if process.returncode != 0 and not script_text:
                stderr_out = process.stderr.read()
                yield f"❌ LLM worker failed (code {process.returncode}): {stderr_out[:500]}", ""
                return

            if not script_text.strip().endswith("}"):
                script_text = script_text.strip() + "\n}"

            yield "✅ Script generation complete!", script_text

        finally:
            try:
                os.unlink(config_path)
            except OSError:
                pass

    except Exception as e:
        yield f"❌ LLM Error: {e}", ""

def build_database(pdf_file):
    if not pdf_file:
        yield "❌ Error: Please upload a PDF in Step 1 first!"
        return
        
    yield "📂 Extracting images and text clusters from PDF..."
    try:
        extract_universal_assets(pdf_file.name, METADATA_FOLDER)
    except Exception as e:
        yield f"❌ PDF Extraction Error: {str(e)}"
        return

    if not os.path.exists(METADATA_PATH):
        yield "❌ Error: metadata.json was not generated!"
        return

    device = "cuda" if torch.cuda.is_available() else "cpu"
    yield f"🔄 Booting up JinaCLIP on {device}..."

    try:
        model = AutoModel.from_pretrained('jinaai/jina-clip-v1', trust_remote_code=True, low_cpu_mem_usage=False, _fast_init=False).to(device)
        client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
        
        try: client.delete_collection(name="course_assets")
        except ValueError: pass
            
        collection = client.create_collection(name="course_assets", metadata={"hnsw:space": "cosine"})
        
        with open(METADATA_PATH, "r", encoding="utf-8") as f:
            metadata = json.load(f)
            
        yield f"🧠 Generating multimodal embeddings for {len(metadata)} images (Batched Processing)..."
        embeddings, metadatas, ids, documents = [], [], [], []
        
        BATCH_SIZE = 16 
        for i in range(0, len(metadata), BATCH_SIZE):
            batch_items = metadata[i : i + BATCH_SIZE]
            valid_images = []
            valid_indices = []
            
            for j, item in enumerate(batch_items):
                if os.path.exists(item["image_path"]):
                    valid_images.append(Image.open(item["image_path"]).convert("RGB"))
                    valid_indices.append(j)
            
            if not valid_images:
                continue
                
            # Process all images in the batch at once ⚡
            batch_embs = model.encode_image(valid_images)

            # Batch-encode captions for hybrid embeddings
            batch_captions = [batch_items[j].get("caption", "") for j in valid_indices]
            caption_emb_map = {}
            non_empty = [(idx, c) for idx, c in enumerate(batch_captions) if c]
            if non_empty:
                cap_texts = [c for _, c in non_empty]
                cap_embs = model.encode_text(cap_texts)
                for pos, (idx, _) in enumerate(non_empty):
                    caption_emb_map[idx] = cap_embs[pos]

            for k, original_j in enumerate(valid_indices):
                item = batch_items[original_j]
                img_emb = batch_embs[k]
                caption = item.get("caption", "")

                # Hybrid embedding: blend image + caption text for better retrieval
                if k in caption_emb_map:
                    txt_emb = caption_emb_map[k]
                    combined = 0.6 * img_emb + 0.4 * txt_emb
                    norm = float((combined ** 2).sum() ** 0.5)
                    if norm > 0:
                        combined = combined / norm
                    embeddings.append(combined.tolist())
                else:
                    embeddings.append(img_emb.tolist())

                metadatas.append({
                    "image_path": item["image_path"],
                    "slide_text": item.get("surrounding_text", "")[:500],
                    "caption": caption
                })
                documents.append((caption if caption else f"page {item.get('page_number', '?')}").lower())
                ids.append(f"asset_{i + original_j}")
                
            yield f"🧠 Processing batch... ({min(i + BATCH_SIZE, len(metadata))}/{len(metadata)} images done)"
            
        yield "💾 Writing embeddings and metadata to ChromaDB..."
        collection.add(embeddings=embeddings, metadatas=metadatas, ids=ids, documents=documents)
        
        del model
        gc.collect()
        torch.cuda.empty_cache()
            
        yield f"✅ Database built successfully! {collection.count()} visual assets stored."
    except Exception as e:
        yield f"❌ Database Error: {str(e)}"

def add_custom_assets(uploaded_files, captions_text):
    """Add user-uploaded images to the existing ChromaDB visual database.
    Each uploaded image gets its own embedding and can be retrieved by fact-grounded queries."""
    if not uploaded_files:
        return "❌ No files uploaded. Please select images first."

    # Parse captions (one per line, matching upload order)
    captions = [c.strip() for c in captions_text.strip().split('\n')] if captions_text and captions_text.strip() else []

    device = "cuda" if torch.cuda.is_available() else "cpu"
    try:
        model = AutoModel.from_pretrained('jinaai/jina-clip-v1', trust_remote_code=True, low_cpu_mem_usage=False, _fast_init=False).to(device)
        client = chromadb.PersistentClient(path=CHROMA_DB_PATH)

        try:
            collection = client.get_collection(name="course_assets")
        except Exception:
            # No database yet — create one
            collection = client.create_collection(name="course_assets", metadata={"hnsw:space": "cosine"})

        existing_count = collection.count()
        os.makedirs(METADATA_FOLDER, exist_ok=True)

        # Load existing metadata.json (or create empty list)
        if os.path.exists(METADATA_PATH):
            with open(METADATA_PATH, "r", encoding="utf-8") as f:
                metadata_log = json.load(f)
        else:
            metadata_log = []

        embeddings, metadatas, ids, documents = [], [], [], []
        added = 0

        for i, upload in enumerate(uploaded_files):
            # Handle both filepath strings and Gradio NamedString/UploadFile objects
            file_path = upload.name if hasattr(upload, 'name') else str(upload)
            if not os.path.exists(file_path):
                continue

            caption = captions[i] if i < len(captions) else ""

            # Copy to assets folder with unique name
            ext = os.path.splitext(file_path)[1] or ".png"
            dest_name = f"custom_asset_{existing_count + added}{ext}"
            dest_path = os.path.join(METADATA_FOLDER, dest_name)
            import shutil
            shutil.copy2(file_path, dest_path)

            # Generate embedding
            img = Image.open(dest_path).convert("RGB")
            img_emb = model.encode_image([img])[0]

            if caption:
                txt_emb = model.encode_text([caption])[0]
                combined = 0.6 * img_emb + 0.4 * txt_emb
                norm = float((combined ** 2).sum() ** 0.5)
                if norm > 0:
                    combined = combined / norm
                embeddings.append(combined.tolist())
            else:
                embeddings.append(img_emb.tolist())

            metadatas.append({
                "image_path": dest_path,
                "slide_text": caption[:500],
                "caption": caption
            })
            documents.append((caption if caption else f"custom asset {added + 1}").lower())
            ids.append(f"custom_{existing_count + added}")

            metadata_log.append({
                "image_path": dest_path,
                "page_number": -1,
                "surrounding_text": caption,
                "caption": caption
            })
            added += 1

        if embeddings:
            collection.add(embeddings=embeddings, metadatas=metadatas, ids=ids, documents=documents)

        # Update metadata.json
        with open(METADATA_PATH, "w", encoding="utf-8") as f:
            json.dump(metadata_log, f, indent=4)

        del model
        gc.collect()
        torch.cuda.empty_cache()

        return f"✅ Added {added} custom asset(s). Database now has {collection.count()} total assets."
    except Exception as e:
        return f"❌ Error adding assets: {str(e)}"

def test_rag_query(query_text):
    if not query_text or str(query_text).strip() == "":
        return None

    # ── Direct slide/image number lookup (bypasses CLIP) ──
    direct_hit = _direct_slide_lookup(query_text)
    if direct_hit:
        return [(direct_hit, "Match #1 | Direct slide/image lookup")]

    device = "cuda" if torch.cuda.is_available() else "cpu"
    try:
        model = AutoModel.from_pretrained('jinaai/jina-clip-v1', trust_remote_code=True, low_cpu_mem_usage=False, _fast_init=False).to(device)
        client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
        collection = client.get_collection(name="course_assets")
        
        text_emb = model.encode_text([query_text])
        results = collection.query(query_embeddings=text_emb.tolist(), n_results=2)

        # If query has a figure/table label, try caption-filtered search first
        label_match = re.search(r'((?:Figure|Fig\.?|Table|Tab\.?|Chart|Diagram|Graph|Equation|Eq\.?)\s*\.?\s*[\(]?\d+[a-z]?[\)]?)', query_text, re.IGNORECASE)
        if label_match:
            for variant in _label_search_variants(label_match.group(1)):
                try:
                    label_results = collection.query(
                        query_embeddings=text_emb.tolist(), n_results=2,
                        where_document={"$contains": variant}
                    )
                    if label_results['metadatas'] and label_results['metadatas'][0]:
                        results = label_results
                        break
                except Exception:
                    continue

        retrieved_gallery = []
        for i in range(len(results['metadatas'][0])):
            img_path = results['metadatas'][0][i]['image_path']
            distance = results['distances'][0][i]
            caption = results['metadatas'][0][i].get('caption', '')
            label_tag = f" | {caption[:60]}" if caption else ""
            retrieved_gallery.append((img_path, f"Match #{i+1} | Distance: {distance:.3f}{label_tag}"))

        del model
        gc.collect()
        torch.cuda.empty_cache()
        return retrieved_gallery
    except Exception as e:
        return [(PLACEHOLDER_IMAGE, f"Error: {e}")]

def step2_fetch_assets(script_json):
    if not script_json or str(script_json).strip() == "":
        return [(PLACEHOLDER_IMAGE, "⚠️ Error: Please generate a JSON script in Tab 1 first!")]

    try:
        script_data = json.loads(script_json)
        if isinstance(script_data, dict) and "segments" in script_data:
            script_data = script_data["segments"]
            
        # Build (segment_index, query) pairs to label gallery items correctly
        query_segments = []
        for idx, item in enumerate(script_data):
            if item.get('visual_mode') == 'RETRIEVE_ASSET' and item.get('query'):
                query_segments.append((idx, item.get('query')))
        queries = [q for _, q in query_segments]
    except Exception as e:
        return [(PLACEHOLDER_IMAGE, f"JSON Parse Error: {e}")]
    
    if not queries: return [(PLACEHOLDER_IMAGE, "No 'RETRIEVE_ASSET' tags found in script.")]

    device = "cuda" if torch.cuda.is_available() else "cpu"
    
    model = AutoModel.from_pretrained('jinaai/jina-clip-v1', trust_remote_code=True, low_cpu_mem_usage=False, _fast_init=False).to(device)
    client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
    
    try: collection = client.get_collection(name="course_assets")
    except (ValueError, Exception): return [(PLACEHOLDER_IMAGE, "Error: Database not built yet.")]
    
    retrieved_gallery = []
    
    for qi, q in enumerate(queries):
        seg_idx = query_segments[qi][0] if qi < len(query_segments) else qi
        query_short = q[:50] + "..." if len(q) > 50 else q
        
        # ── Direct slide/image number lookup (no CLIP needed) ──
        direct_hit = _direct_slide_lookup(q)
        if direct_hit:
            retrieved_gallery.append((direct_hit, f"Segment {seg_idx + 1}: {query_short}"))
            continue
        
        # ── CLIP embedding search (for descriptive queries / paper figures) ──
        text_emb = model.encode_text([q])
        results = collection.query(query_embeddings=text_emb.tolist(), n_results=1)

        # If query references a specific figure/table label, try caption-filtered search
        label_match = re.search(r'((?:Figure|Fig\.?|Table|Tab\.?|Chart|Diagram|Graph|Equation|Eq\.?)\s*\.?\s*[\(]?\d+[a-z]?[\)]?)', q, re.IGNORECASE)
        if label_match:
            for variant in _label_search_variants(label_match.group(1)):
                try:
                    label_results = collection.query(
                        query_embeddings=text_emb.tolist(),
                        n_results=1,
                        where_document={"$contains": variant}
                    )
                    if label_results['metadatas'] and label_results['metadatas'][0]:
                        results = label_results
                        break
                except Exception:
                    continue

        if results['metadatas'] and results['metadatas'][0]:
            best_img_path = results['metadatas'][0][0]['image_path']
            distance = results['distances'][0][0]
            retrieved_gallery.append((best_img_path, f"Segment {seg_idx + 1}: {query_short}"))

    del model
    gc.collect()
    torch.cuda.empty_cache()
    return retrieved_gallery

# Avatar generation helper (registry-driven; default backend = SoulX-FlashHead)
def generate_talking_head_avatar(
    audio_path: str,
    source_image_path: str = None,
    output_dir: str = "final_render_assets",
    progress_callback=None,
    avatar_root: str = "",
    avatar_env_python: str = "",
    avatar_backend: str = "",
) -> str | None:
    """Avatar generation — dispatches to the chosen backend in avatar_backends.py.

    `avatar_backend` may be either a registry key (e.g. "soulx-flashhead") or
    the user-facing dropdown label. Defaults to DEFAULT_AVATAR_BACKEND.
    """

    # Resolve backend, paths
    backend = get_avatar_backend(avatar_backend or DEFAULT_AVATAR_BACKEND)
    active_root = avatar_root.strip() if avatar_root and avatar_root.strip() else DEFAULT_AVATAR_ROOT
    active_python = avatar_env_python.strip() if avatar_env_python and avatar_env_python.strip() else DEFAULT_AVATAR_ENV_PYTHON

    if not os.path.isdir(active_root):
        print(f"❌ Avatar model root not found: {active_root}")
        return None
    if not os.path.exists(active_python):
        print(f"❌ Avatar venv python not found: {active_python}")
        return None

    if not os.path.exists(audio_path):
        print(f"❌ Audio not found: {audio_path}")
        return None

    audio_path = os.path.abspath(audio_path)
    output_dir = os.path.abspath(output_dir)
    os.makedirs(output_dir, exist_ok=True)

    # Use uploaded headshot, or fall back to a sample shipped with the model repo
    if not source_image_path or not os.path.exists(source_image_path):
        source_image_path = os.path.join(active_root, "examples", "source_image", "man.png")
    source_image_path = os.path.abspath(source_image_path)

    output_video = os.path.join(output_dir, f"avatar_{int(time.time())}.mp4")

    # Auto-discover model artifacts inside the root directory.
    # Convention: ckpt_dir = first non-encoder dir under <root>/models/,
    # audio encoder dir = first dir whose name contains the backend's hint string.
    models_dir = os.path.join(active_root, "models")
    ckpt_dir = None
    audio_encoder_dir = None
    if os.path.isdir(models_dir):
        for d in sorted(os.listdir(models_dir)):
            dp = os.path.join(models_dir, d)
            if not os.path.isdir(dp):
                continue
            if backend.audio_encoder_dir_hint and backend.audio_encoder_dir_hint.lower() in d.lower():
                audio_encoder_dir = dp
            elif ckpt_dir is None:
                ckpt_dir = dp  # first non-encoder dir is the main checkpoint

    # Resolve inference script inside the repo root
    inference_script = os.path.join(active_root, backend.entry_script)
    if not os.path.exists(inference_script):
        print(f"❌ Inference script not found for backend '{backend.name}': {inference_script}")
        return None

    cmd = build_avatar_command(
        backend=backend,
        python_exe=active_python,
        repo_root=active_root,
        inference_script_path=inference_script,
        image_path=source_image_path,
        audio_path=audio_path,
        output_path=output_video,
        ckpt_dir=ckpt_dir,
        audio_encoder_dir=audio_encoder_dir,
    )

    print(f"🚀 Running avatar backend: {backend.name}")
    print(f"   Repo: {active_root}")
    print(f"   Python: {active_python}")
    if ckpt_dir:
        print(f"   Checkpoint: {ckpt_dir}")
    if audio_encoder_dir:
        print(f"   Audio encoder: {audio_encoder_dir}")
    print(f"   Image: {os.path.basename(source_image_path)}")
    print(f"   Audio: {os.path.basename(audio_path)}")

    try:
        # Disable torch.compile to avoid Triton incompatibility on Windows
        env = os.environ.copy()
        env["TORCH_COMPILE_DISABLE"] = "1"
        # Sanitize PYTHONHASHSEED — upstream libs may set it > 4294967295
        _hs = env.get("PYTHONHASHSEED", "")
        if _hs and _hs != "random":
            try:
                if int(_hs) > 4294967295:
                    env["PYTHONHASHSEED"] = str(int(_hs) % 4294967296)
            except ValueError:
                env.pop("PYTHONHASHSEED", None)

        process = subprocess.Popen(
            cmd,
            cwd=active_root,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            bufsize=1,
            universal_newlines=True,
            env=env
        )

        # Forward progress if the script prints percentages
        for line in process.stdout:
            if progress_callback and ("%" in line or "progress" in line.lower()):
                progress_callback(0.5, f"{backend.name} rendering... {line.strip()}")
            print(line.strip())

        process.wait()

        # Most backends write directly to `output_video`. A few (e.g. SadTalker)
        # take a result *directory* and write a timestamped file inside it — in
        # that case, pick the newest .mp4 produced after the run started.
        produced_path = output_video
        if process.returncode == 0 and not os.path.exists(produced_path):
            search_dir = output_dir
            if os.path.isdir(search_dir):
                mp4s = [os.path.join(search_dir, f) for f in os.listdir(search_dir)
                        if f.lower().endswith(".mp4")]
                if mp4s:
                    produced_path = max(mp4s, key=os.path.getmtime)

        if process.returncode == 0 and os.path.exists(produced_path):
            print(f"✅ Avatar success → {produced_path}")
            return produced_path
        else:
            print(f"❌ Avatar backend '{backend.name}' failed (code {process.returncode})")
            return None

    except Exception as e:
        print(f"❌ Avatar Error: {e}")
        return None
        
def create_static_video_clip(image_path, audio_path, output_path, canvas_w=512, canvas_h=512):
    """Converts a static image and audio file into a single MP4 clip at specified canvas resolution.
    Uses audio duration to determine video length (no -shortest flag), preventing audio truncation."""
    # Get audio duration to set explicit video length (avoids -shortest truncation)
    probe_cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", audio_path]
    try:
        duration = float(subprocess.run(probe_cmd, capture_output=True, text=True, check=True).stdout.strip())
    except Exception:
        duration = None

    scale_filter = (
        f"scale={canvas_w}:{canvas_h}:force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={canvas_w}:{canvas_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1"
    )
    cmd = [
        "ffmpeg", "-y", "-loop", "1", "-i", image_path, "-i", audio_path,
        "-vf", scale_filter,
        "-c:v", "libx264", "-preset", "fast", "-crf", "23",
        "-c:a", "aac", "-b:a", "192k",
        "-r", "25", "-pix_fmt", "yuv420p",
    ]
    if duration is not None:
        cmd.extend(["-t", f"{duration:.3f}"])
    else:
        cmd.append("-shortest")  # Fallback if probe fails
    cmd.append(output_path)
    subprocess.run(cmd, check=True, capture_output=True)
    return output_path

def finalize_full_video(output_folder, timeline, canvas_w=1280, canvas_h=720, enable_subtitles=True):
    """Stitches all segment MP4s into one FINAL_LECTURE.mp4 at consistent canvas resolution,
    with optionally burned-in subtitles generated from the spoken_text timeline."""
    concat_file = os.path.join(output_folder, "concat_list.txt")
    srt_file = os.path.join(output_folder, "subtitles.srt")
    final_output = os.path.join(output_folder, "FINAL_LECTURE.mp4")

    # --- Generate SRT subtitle file from timeline ---
    cumulative_time = 0.0
    srt_entries = []
    for entry_idx, entry in enumerate(timeline):
        spoken_text = entry.get("spoken_text", "") or entry.get("script", "") or ""
        spoken_text = spoken_text.strip()
        if not spoken_text:
            spoken_text = entry.get("query", "").strip()

        # Clean TTS breathing artifacts from subtitle text
        import re as _re
        spoken_text = _latex_to_subtitle(spoken_text)  # Convert LaTeX to clean display text
        spoken_text = _re.sub(r'\.{2,}', '.', spoken_text)   # "..." → "."
        spoken_text = _re.sub(r'\.\s*\.', '.', spoken_text)   # ". ." → "."
        spoken_text = _re.sub(r'\.{2,}', '.', spoken_text)    # catch leftovers
        # Clean em/en dashes → commas, fix word concatenation
        spoken_text = spoken_text.replace('\u2014', ',').replace('\u2013', ',')
        spoken_text = _re.sub(r'([a-z])([A-Z])', r'\1 \2', spoken_text)
        spoken_text = _re.sub(r'([.!?,;:])([A-Za-z])', r'\1 \2', spoken_text)
        spoken_text = _re.sub(r'  +', ' ', spoken_text)
        spoken_text = spoken_text.strip()

        # Get segment duration from audio file
        audio_path = os.path.join(output_folder, entry.get("audio_file", ""))
        seg_duration = 5.0  # default
        if os.path.exists(audio_path):
            try:
                probe_cmd = [
                    "ffprobe", "-v", "quiet", "-show_entries", "format=duration",
                    "-of", "default=noprint_wrappers=1:nokey=1", audio_path
                ]
                result = subprocess.run(probe_cmd, capture_output=True, text=True)
                seg_duration = float(result.stdout.strip())
            except Exception:
                pass

        start_time = cumulative_time
        end_time = cumulative_time + seg_duration
        cumulative_time = end_time

        # Split long text into ~80 char subtitle chunks
        words = spoken_text.split()
        chunks = []
        current = []
        current_len = 0
        for w in words:
            if current_len + len(w) + 1 > 80 and current:
                chunks.append(" ".join(current))
                current = [w]
                current_len = len(w)
            else:
                current.append(w)
                current_len += len(w) + 1
        if current:
            chunks.append(" ".join(current))

        if not chunks:
            chunks = [spoken_text or "..."]

        chunk_duration = (end_time - start_time) / len(chunks) if chunks else (end_time - start_time)
        for ci, chunk_text in enumerate(chunks):
            c_start = start_time + ci * chunk_duration
            c_end = start_time + (ci + 1) * chunk_duration
            srt_entries.append((len(srt_entries) + 1, c_start, c_end, chunk_text))

    def _fmt_srt_time(seconds):
        h = int(seconds // 3600)
        m = int((seconds % 3600) // 60)
        s = int(seconds % 60)
        ms = int((seconds - int(seconds)) * 1000)
        return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"

    with open(srt_file, "w", encoding="utf-8") as f:
        for idx, start, end, text in srt_entries:
            f.write(f"{idx}\n{_fmt_srt_time(start)} --> {_fmt_srt_time(end)}\n{text}\n\n")

    # --- Build concat list ---
    with open(concat_file, "w") as f:
        for entry in timeline:
            if entry.get("video_file") and os.path.exists(os.path.join(output_folder, entry["video_file"])):
                clip_path = entry["video_file"]
            else:
                img_path = os.path.join(output_folder, entry["image_file"])
                aud_path = os.path.join(output_folder, entry["audio_file"])
                clip_name = f"temp_clip_{entry['segment']}.mp4"
                clip_path = create_static_video_clip(
                    img_path, aud_path, 
                    os.path.join(output_folder, clip_name),
                    canvas_w=canvas_w, canvas_h=canvas_h
                )
            f.write(f"file '{os.path.basename(clip_path)}'\n")

    # --- Concat + optionally burn subtitles ---
    scale_filter = (
        f"scale={canvas_w}:{canvas_h}:force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={canvas_w}:{canvas_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1"
    )
    if enable_subtitles:
        # Escape backslashes and colons in the SRT path for FFmpeg filter
        srt_escaped = srt_file.replace("\\", "/").replace(":", "\\:")
        scale_filter += (
            f",subtitles='{srt_escaped}':force_style='FontSize=13,PrimaryColour=&H00FFFFFF,OutlineColour=&H00000000,Outline=2,Shadow=1,MarginV=15'"
        )
    cmd = [
        "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_file,
        "-vf", scale_filter,
        "-c:v", "libx264", "-preset", "fast", "-crf", "23",
        "-c:a", "aac", "-b:a", "192k",
        "-r", "25", "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        final_output
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    return final_output

# ==========================================
# 🎬 PIP COMPOSITING & RESOLUTION SYSTEM
# ==========================================
def fetch_pexels_broll(query, orientation="landscape", min_duration=5, per_page=3, return_multiple=False):
    """Download matching B-roll video(s) from Pexels and cache locally.
    
    Returns the local path to the downloaded video, or None on failure.
    If return_multiple=True, returns a list of local paths (for CLIP reranking).
    Downloaded files are saved to BROLL_ASSETS_DIR for reuse.
    
    Args:
        query: Search terms (e.g. "city traffic night")
        orientation: landscape | portrait | square
        min_duration: Minimum video duration in seconds
        per_page: Number of results to fetch
        return_multiple: If True, download and return up to per_page candidates
    """
    if not PEXELS_API_KEY:
        print("⚠️ Pexels API key not set. Set PEXELS_API_KEY env var or paste in app.py.")
        return [] if return_multiple else None
    
    try:
        headers = {"Authorization": PEXELS_API_KEY}
        params = {
            "query": query,
            "orientation": orientation,
            "per_page": max(per_page, 5) if return_multiple else per_page,
            "size": "medium",
        }
        
        resp = requests.get(
            "https://api.pexels.com/videos/search",
            headers=headers,
            params=params,
            timeout=15
        )
        resp.raise_for_status()
        data = resp.json()
        
        videos = data.get("videos", [])
        if not videos:
            print(f"🔍 Pexels: No videos found for '{query}'")
            return [] if return_multiple else None
        
        downloaded_paths = []
        
        for video in videos:
            if video.get("duration", 0) < min_duration:
                continue
            
            # Sort video files by resolution (height), pick highest ≤ 1080p
            files = sorted(
                video.get("video_files", []),
                key=lambda f: f.get("height", 0),
                reverse=True
            )
            
            download_url = None
            for vf in files:
                h = vf.get("height", 0)
                if 360 <= h <= 1080 and vf.get("link"):
                    download_url = vf["link"]
                    break
            
            if not download_url and files:
                download_url = files[0].get("link")
            
            if download_url:
                safe_name = re.sub(r'[^a-zA-Z0-9_]', '_', query.lower()).strip('_')[:50]
                local_filename = f"pexels_{safe_name}_{video['id']}.mp4"
                local_path = os.path.join(BROLL_ASSETS_DIR, local_filename)
                
                # Skip download if already cached
                if os.path.exists(local_path) and os.path.getsize(local_path) > 10000:
                    print(f"✅ Pexels cache hit: {local_filename}")
                    if return_multiple:
                        downloaded_paths.append(local_path)
                        continue
                    else:
                        return local_path
                
                print(f"⬇️  Pexels: Downloading B-roll for '{query}' → {local_filename}")
                dl_resp = requests.get(download_url, stream=True, timeout=60)
                dl_resp.raise_for_status()
                with open(local_path, "wb") as f:
                    for chunk in dl_resp.iter_content(chunk_size=1024 * 256):
                        f.write(chunk)
                
                if os.path.exists(local_path) and os.path.getsize(local_path) > 10000:
                    print(f"✅ Pexels: Downloaded {local_filename} ({os.path.getsize(local_path) / 1024:.0f} KB)")
                    if return_multiple:
                        downloaded_paths.append(local_path)
                    else:
                        return local_path
                else:
                    print(f"⚠️ Pexels: Download too small or failed for '{query}'")
                    if os.path.exists(local_path):
                        os.remove(local_path)
        
        if return_multiple and downloaded_paths:
            return downloaded_paths
        
        if not return_multiple and not downloaded_paths:
            print(f"🔍 Pexels: No suitable video found for '{query}'")
        return downloaded_paths if return_multiple else None
        
    except requests.exceptions.RequestException as e:
        print(f"⚠️ Pexels API error: {e}")
        return [] if return_multiple else None
    except Exception as e:
        print(f"⚠️ Pexels unexpected error: {e}")
        return [] if return_multiple else None

def _extract_middle_frame(video_path):
    """Extract a single frame from the middle of a video file using FFmpeg.
    Returns a PIL Image or None on failure."""
    try:
        # Get duration
        probe_cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration",
                     "-of", "default=noprint_wrappers=1:nokey=1", video_path]
        result = subprocess.run(probe_cmd, capture_output=True, text=True, timeout=10)
        duration = float(result.stdout.strip()) if result.stdout.strip() else 5.0
        mid_time = duration / 2.0
        
        # Extract frame at midpoint
        temp_frame = os.path.join(BROLL_ASSETS_DIR, f"_temp_frame_{os.getpid()}.jpg")
        cmd = ["ffmpeg", "-y", "-ss", f"{mid_time:.2f}", "-i", video_path,
               "-frames:v", "1", "-q:v", "2", temp_frame]
        subprocess.run(cmd, capture_output=True, timeout=15)
        
        if os.path.exists(temp_frame) and os.path.getsize(temp_frame) > 100:
            img = Image.open(temp_frame).convert("RGB")
            os.remove(temp_frame)
            return img
    except Exception as e:
        print(f"⚠️ Frame extraction failed for {video_path}: {e}")
    return None


def rerank_broll_with_clip(candidates, narration_text, visual_model):
    """Rerank B-roll video candidates using JinaCLIP text-image similarity.
    
    Extracts a middle frame from each candidate video, embeds it alongside
    the narration text, and returns the candidate with highest cosine similarity.
    
    Args:
        candidates: List of local video file paths
        narration_text: The spoken_text for this segment (used as semantic query)
        visual_model: Loaded JinaCLIP model instance
    
    Returns:
        Best matching video path, or candidates[0] as fallback
    """
    if not candidates:
        return None
    if len(candidates) == 1 or visual_model is None:
        return candidates[0]
    
    try:
        # Extract middle frames from each candidate
        frames = []
        valid_candidates = []
        for path in candidates:
            frame = _extract_middle_frame(path)
            if frame is not None:
                frames.append(frame)
                valid_candidates.append(path)
        
        if not frames:
            return candidates[0]
        
        # Encode all frames and the narration text
        img_embs = visual_model.encode_image(frames)
        text_emb = visual_model.encode_text([narration_text])[0]
        
        # Cosine similarity (embeddings are already normalized by JinaCLIP)
        similarities = []
        for img_emb in img_embs:
            sim = float(np.dot(text_emb, img_emb) / (np.linalg.norm(text_emb) * np.linalg.norm(img_emb) + 1e-8))
            similarities.append(sim)
        
        best_idx = int(np.argmax(similarities))
        best_path = valid_candidates[best_idx]
        print(f"🎯 CLIP rerank: selected '{os.path.basename(best_path)}' "
              f"(sim={similarities[best_idx]:.3f}) from {len(valid_candidates)} candidates")
        
        for i, (path, sim) in enumerate(zip(valid_candidates, similarities)):
            marker = " ← selected" if i == best_idx else ""
            print(f"    [{i+1}] {os.path.basename(path)}: {sim:.3f}{marker}")
        
        return best_path
        
    except Exception as e:
        print(f"⚠️ CLIP reranking failed, using first candidate: {e}")
        return candidates[0]


def find_broll_asset(query, broll_query=None, narration_text=None, visual_model=None):
    """Find matching B-roll from broll_assets/ folder using keyword matching,
    with optional CLIP-based reranking of Pexels results.
    
    Args:
        query: Topic-level search query (from segment 'query' field)
        broll_query: Visually descriptive search phrase (from segment 'broll_query' field)
        narration_text: Spoken text for this segment (used for CLIP reranking)
        visual_model: Loaded JinaCLIP model for CLIP reranking (optional)
    """
    search_query = broll_query if broll_query else query
    
    if not os.path.exists(BROLL_ASSETS_DIR):
        os.makedirs(BROLL_ASSETS_DIR, exist_ok=True)
    
    broll_files = []
    for f in os.listdir(BROLL_ASSETS_DIR):
        ext = os.path.splitext(f)[1].lower()
        if ext in {'.mp4', '.mov', '.avi', '.webm', '.mkv', '.jpg', '.jpeg', '.png', '.gif'}:
            broll_files.append(f)
    
    # Keyword matching against local files
    if broll_files:
        query_words = set(re.sub(r'[_\-./\\]', ' ', search_query.lower()).split())
        best_match = None
        best_score = 0
        
        for f in broll_files:
            name_words = set(re.sub(r'[_\-./\\]', ' ', os.path.splitext(f)[0].lower()).split())
            score = len(query_words & name_words)
            if score > best_score:
                best_score = score
                best_match = f
        
        if best_match and best_score > 0:
            return os.path.join(BROLL_ASSETS_DIR, best_match)
    
    # No local match → fetch from Pexels with CLIP reranking
    use_reranking = visual_model is not None and narration_text
    
    if use_reranking:
        # Fetch multiple candidates for CLIP reranking
        candidates = fetch_pexels_broll(search_query, per_page=5, return_multiple=True)
        if candidates:
            return rerank_broll_with_clip(candidates, narration_text, visual_model)
    
    # Fallback: single Pexels result (no CLIP model available)
    pexels_result = fetch_pexels_broll(search_query)
    if pexels_result:
        return pexels_result
    
    # Nothing found at all
    return None

def create_pip_composite(bg_path, avatar_video, audio_path, output_path,
                         canvas_w=1280, canvas_h=720, pip_ratio=0.25, 
                         pip_padding=20, bg_is_video=False):
    """
    Universal PiP compositor using FFmpeg.
    
    Creates a video with:
    - Background: high-res image or video (scaled to canvas with Lanczos)
    - Overlay: talking-head avatar in upper-right corner (~25% size)
    - Audio: from the TTS pipeline
    
    This decouples avatar resolution (512x512) from asset resolution (native).
    """
    pip_w = int(canvas_w * pip_ratio)
    pip_h = int(canvas_h * pip_ratio)
    overlay_x = canvas_w - pip_w - pip_padding
    overlay_y = pip_padding
    
    # Build filter complex
    bg_scale = (
        f"[0:v]scale={canvas_w}:{canvas_h}:force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={canvas_w}:{canvas_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1,format=yuv420p[bg]"
    )
    # tpad=stop_mode=clone freezes avatar on last frame so audio is never truncated
    avatar_scale = (
        f"[1:v]scale={pip_w}:{pip_h}:force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={pip_w}:{pip_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1,format=yuv420p,"
        f"tpad=stop_mode=clone:stop=-1[pip]"
    )
    overlay = f"[bg][pip]overlay={overlay_x}:{overlay_y}:eof_action=repeat[out]"
    filter_complex = f"{bg_scale};{avatar_scale};{overlay}"

    # Get audio duration to set explicit output length (avoids -shortest truncation)
    probe_cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", audio_path]
    try:
        audio_duration = float(subprocess.run(probe_cmd, capture_output=True, text=True, check=True).stdout.strip())
    except Exception:
        audio_duration = None

    if bg_is_video:
        cmd = [
            "ffmpeg", "-y",
            "-stream_loop", "-1",   # Loop B-roll if shorter than audio
            "-i", bg_path,          # Background video
            "-i", avatar_video,     # Avatar overlay
            "-i", audio_path,       # TTS audio (overrides B-roll audio)
            "-filter_complex", filter_complex,
            "-map", "[out]", "-map", "2:a",
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-c:a", "aac", "-b:a", "192k",
            "-r", "25", "-pix_fmt", "yuv420p",
        ]
    else:
        cmd = [
            "ffmpeg", "-y",
            "-loop", "1", "-i", bg_path,   # Background image (looped)
            "-i", avatar_video,              # Avatar overlay
            "-i", audio_path,                # TTS audio
            "-filter_complex", filter_complex,
            "-map", "[out]", "-map", "2:a",
            "-c:v", "libx264", "-preset", "fast", "-crf", "23",
            "-c:a", "aac", "-b:a", "192k",
            "-r", "25", "-pix_fmt", "yuv420p",
        ]
    if audio_duration is not None:
        cmd.extend(["-t", f"{audio_duration:.3f}"])
    else:
        cmd.append("-shortest")  # Fallback if probe fails
    cmd.append(output_path)
    
    print(f"🎬 PiP composite: {os.path.basename(bg_path)} + avatar → {os.path.basename(output_path)}")
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"⚠️ FFmpeg PiP error: {result.stderr[:500]}")
        return None
    return output_path

def upscale_avatar_video(avatar_video, output_path, canvas_w=1280, canvas_h=720):
    """Upscale standalone avatar video to canvas resolution with Lanczos + black padding."""
    scale_filter = (
        f"scale={canvas_w}:{canvas_h}:force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={canvas_w}:{canvas_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1"
    )
    cmd = [
        "ffmpeg", "-y", "-i", avatar_video,
        "-vf", scale_filter,
        "-c:v", "libx264", "-preset", "fast", "-crf", "23",
        "-c:a", "copy",
        "-r", "25", "-pix_fmt", "yuv420p", output_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"⚠️ Upscale error: {result.stderr[:500]}")
        return None
    return output_path
    
# ==========================================
# LATEX MATH → SPOKEN / DISPLAY TEXT CONVERTERS
# ==========================================
# Greek letter map (LaTeX command → spoken name)
_GREEK_MAP = {
    'alpha': 'alpha', 'beta': 'beta', 'gamma': 'gamma', 'delta': 'delta',
    'epsilon': 'epsilon', 'zeta': 'zeta', 'eta': 'eta', 'theta': 'theta',
    'iota': 'iota', 'kappa': 'kappa', 'lambda': 'lambda', 'mu': 'mu',
    'nu': 'nu', 'xi': 'xi', 'pi': 'pi', 'rho': 'rho',
    'sigma': 'sigma', 'tau': 'tau', 'upsilon': 'upsilon', 'phi': 'phi',
    'chi': 'chi', 'psi': 'psi', 'omega': 'omega',
    'Alpha': 'Alpha', 'Beta': 'Beta', 'Gamma': 'Gamma', 'Delta': 'Delta',
    'Theta': 'Theta', 'Lambda': 'Lambda', 'Sigma': 'Sigma', 'Omega': 'Omega',
    'Phi': 'Phi', 'Psi': 'Psi', 'Pi': 'Pi',
    'nabla': 'nabla', 'partial': 'partial', 'infty': 'infinity',
    'cdot': 'dot', 'cdots': 'dots', 'ldots': 'dots', 'times': 'times',
    'approx': 'approximately', 'neq': 'not equal to', 'leq': 'less than or equal to',
    'geq': 'greater than or equal to', 'rightarrow': 'arrow',
    'Rightarrow': 'implies', 'leftarrow': 'arrow', 'sum': 'the sum',
    'prod': 'the product', 'int': 'the integral', 'sqrt': 'the square root of',
    'text': '', 'mathrm': '', 'mathbf': '', 'mathcal': '', 'hat': '',
    'bar': '', 'vec': '', 'dot': '', 'ddot': '', 'tilde': '',
}

# Greek letter map for subtitle display (LaTeX command → Unicode)
_GREEK_DISPLAY = {
    'alpha': 'α', 'beta': 'β', 'gamma': 'γ', 'delta': 'δ',
    'epsilon': 'ε', 'zeta': 'ζ', 'eta': 'η', 'theta': 'θ',
    'iota': 'ι', 'kappa': 'κ', 'lambda': 'λ', 'mu': 'μ',
    'nu': 'ν', 'xi': 'ξ', 'pi': 'π', 'rho': 'ρ',
    'sigma': 'σ', 'tau': 'τ', 'upsilon': 'υ', 'phi': 'φ',
    'chi': 'χ', 'psi': 'ψ', 'omega': 'ω',
    'Alpha': 'Α', 'Beta': 'Β', 'Gamma': 'Γ', 'Delta': 'Δ',
    'Theta': 'Θ', 'Lambda': 'Λ', 'Sigma': 'Σ', 'Omega': 'Ω',
    'Phi': 'Φ', 'Psi': 'Ψ', 'Pi': 'Π',
    'nabla': '∇', 'partial': '∂', 'infty': '∞',
    'cdot': '·', 'cdots': '⋯', 'ldots': '…', 'times': '×',
    'approx': '≈', 'neq': '≠', 'leq': '≤', 'geq': '≥',
    'rightarrow': '→', 'Rightarrow': '⇒', 'leftarrow': '←',
    'sum': '∑', 'prod': '∏', 'int': '∫', 'sqrt': '√',
    'text': '', 'mathrm': '', 'mathbf': '', 'mathcal': '', 'hat': '',
    'bar': '', 'vec': '', 'dot': '', 'ddot': '', 'tilde': '',
}

# Unicode subscript/superscript character maps
_UNICODE_SUB = str.maketrans('0123456789+-=()aehijklmnoprstuvx',
                              '₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎ₐₑₕᵢⱼₖₗₘₙₒₚᵣₛₜᵤᵥₓ')
_UNICODE_SUP = str.maketrans('0123456789+-=()abcdefghijklmnoprstuvwxyz',
                              '⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ᵃᵇᶜᵈᵉᶠᵍʰⁱʲᵏˡᵐⁿᵒᵖʳˢᵗᵘᵛʷˣʸᶻ')


def _latex_to_spoken(text):
    """Convert LaTeX math notation in text to natural spoken English for TTS.
    Handles $...$ inline math, subscripts, superscripts, fractions, Greek letters."""
    if not text or '$' not in text and '\\' not in text:
        return text

    def _convert_math_block(math_str):
        """Convert a single LaTeX math expression to spoken words."""
        s = math_str.strip()

        # Replace \frac{a}{b} → "a over b"
        while '\\frac' in s:
            m = re.search(r'\\frac\s*\{([^{}]*)\}\s*\{([^{}]*)\}', s)
            if not m:
                break
            num, den = m.group(1).strip(), m.group(2).strip()
            s = s[:m.start()] + f'{num} over {den}' + s[m.end():]

        # Replace Greek letters and LaTeX commands: \beta → "beta"
        for cmd, spoken in sorted(_GREEK_MAP.items(), key=lambda x: -len(x[0])):
            s = s.replace(f'\\{cmd}', f' {spoken} ' if spoken else ' ')

        # Replace subscripts FIRST (resolves inner braces for nested ^{\beta_{jj}})
        while '_{' in s:
            m = re.search(r'_\{([^{}]*)\}', s)
            if not m:
                break
            sub = m.group(1).strip()
            # Handle comma-separated subscripts: "f,jj" → "f, j j"
            if ',' in sub:
                parts = [p.strip() for p in sub.split(',')]
                parts = [' '.join(p) if p.isalpha() and len(p) <= 4 else p for p in parts]
                sub = ', '.join(parts)
            elif sub.isalpha() and len(sub) <= 4:
                sub = ' '.join(sub)
            s = s[:m.start()] + f' sub {sub} ' + s[m.end():]
        # Single-char subscript: _i → "sub i"
        s = re.sub(r'_([0-9a-zA-Z])', r' sub \1', s)

        # Replace superscripts: ^{...} → "to the power ..."
        while '^{' in s:
            m = re.search(r'\^\{([^{}]*)\}', s)
            if not m:
                break
            exp = m.group(1).strip()
            # Handle negative exponents naturally
            if exp.startswith('-'):
                exp = 'negative ' + exp[1:].strip()
            # Convert / to "over" for inline fractions like E/RT
            exp = exp.replace('/', ' over ')
            s = s[:m.start()] + f', to the power {exp}, ' + s[m.end():]
        # Single-char superscript: ^2 → "squared", ^3 → "cubed", else "to the power N"
        def _sup_single(m):
            c = m.group(1)
            if c == '2': return ' squared'
            if c == '3': return ' cubed'
            return f' to the power {c}'
        s = re.sub(r'\^([0-9a-zA-Z])', _sup_single, s)

        # Clean up remaining LaTeX artifacts
        s = s.replace('{', '').replace('}', '').replace('\\', '')
        s = re.sub(r'\s+', ' ', s).strip()
        return s

    # Process $...$ inline math (non-greedy)
    result = re.sub(
        r'\$([^$]+)\$',
        lambda m: ' ' + _convert_math_block(m.group(1)) + ' ',
        text
    )

    # Process $$...$$ display math (if any)
    result = re.sub(
        r'\$\$([^$]+)\$\$',
        lambda m: ' ' + _convert_math_block(m.group(1)) + ' ',
        result
    )

    # Clean double spaces
    result = re.sub(r'\s+', ' ', result)
    # Clean space before punctuation
    result = re.sub(r'\s+([,.\?!;:])', r'\1', result)
    # Clean comma artifacts before punctuation: ",. " → ". "
    result = re.sub(r',\s*([.?!;])', r'\1', result)
    return result.strip()


def _latex_to_subtitle(text):
    """Convert LaTeX math notation to clean Unicode display text for subtitles.
    Uses Unicode subscript/superscript characters and Greek symbols."""
    if not text or '$' not in text and '\\' not in text:
        return text

    def _convert_display(math_str):
        """Convert a single LaTeX math expression to Unicode display."""
        s = math_str.strip()

        # Replace \frac{a}{b} → "a/b"
        while '\\frac' in s:
            m = re.search(r'\\frac\s*\{([^{}]*)\}\s*\{([^{}]*)\}', s)
            if not m:
                break
            num, den = m.group(1).strip(), m.group(2).strip()
            s = s[:m.start()] + f'{num}/{den}' + s[m.end():]

        # Replace Greek letters with Unicode
        for cmd, char in sorted(_GREEK_DISPLAY.items(), key=lambda x: -len(x[0])):
            s = s.replace(f'\\{cmd}', char if char else '')

        # Convert subscripts to Unicode: _{jj} → subscript chars
        while '_{' in s:
            m = re.search(r'_\{([^{}]*)\}', s)
            if not m:
                break
            sub_text = m.group(1).strip().lower()
            sub_unicode = sub_text.translate(_UNICODE_SUB)
            s = s[:m.start()] + sub_unicode + s[m.end():]
        s = re.sub(r'_([0-9a-zA-Z])', lambda m: m.group(1).lower().translate(_UNICODE_SUB), s)

        # Convert superscripts to Unicode: ^{2} → superscript chars
        _sup_capable = set('0123456789+-=()abcdefghijklmnoprstuvwxyz ')
        while '^{' in s:
            m = re.search(r'\^\{([^{}]*)\}', s)
            if not m:
                break
            sup_text = m.group(1).strip()
            sup_lower = sup_text.lower()
            # Check if all chars can be Unicode-superscripted
            if all(c in _sup_capable for c in sup_lower):
                if '-' in sup_lower:
                    sup_unicode = '⁻' + sup_lower.replace('-', '').translate(_UNICODE_SUP)
                else:
                    sup_unicode = sup_lower.translate(_UNICODE_SUP)
            else:
                # Fallback notation for complex superscripts (Greek, subscript chars, etc.)
                sup_unicode = '^(' + sup_text + ')'
            s = s[:m.start()] + sup_unicode + s[m.end():]
        s = re.sub(r'\^([0-9a-zA-Z])', lambda m: m.group(1).lower().translate(_UNICODE_SUP), s)

        # Clean remaining LaTeX
        s = s.replace('{', '').replace('}', '').replace('\\', '')
        s = re.sub(r'\s+', ' ', s).strip()
        return s

    # Process $$...$$ first, then $...$
    result = re.sub(
        r'\$\$([^$]+)\$\$',
        lambda m: _convert_display(m.group(1)),
        text
    )
    result = re.sub(
        r'\$([^$]+)\$',
        lambda m: _convert_display(m.group(1)),
        result
    )
    return result


# ==========================================
# SEMANTIC ENGINE: SMOOTH AUTOREGRESSIVE AUDIO
# ==========================================
# ==========================================
# ADVANCED BREATHING PAUSE SYSTEM FOR MELOTTS
# ==========================================
def add_breathing_pauses(text):
    """
    Enhanced MeloTTS breathing control using strategic punctuation.
    
    MeloTTS (FastSpeech2) duration model respects punctuation for natural pauses:
    - Comma (,) = ~200-300ms pause
    - Semicolon (;) = ~350-400ms pause  
    - Em-dash (—) = ~400-500ms pause
    - Period/Ellipsis = ~500-700ms pause
    
    Strategy: Insert natural breathing points WITHIN sentences, not just at ends.
    This creates natural phrasing and prevents rushed, monotonous speech.
    """
    if not text or not isinstance(text, str):
        return text
    
    text = text.rstrip()
    if len(text) == 0:
        return text
    
    # Split into sentences while preserving punctuation
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    processed_sentences = []
    
    for sentence in sentences:
        if not sentence.strip():
            continue
            
        # For very long sentences (>20 words), add breathing commas at clause boundaries
        words = sentence.split()
        if len(words) > 20:
            # Add comma after relative clauses (after "which", "that", "who", etc.)
            sentence = re.sub(r'\b(which|that|who|whom|whose|where|when)\s+', r'\1 ', sentence)
            
            # Add commas before "and"/"but" when they connect long independent clauses
            sentence = re.sub(r'([,.])\s+(and|but|or)\s+', r'\1, \2 ', sentence)
            
            # Break at major clause boundaries: before coordinating conjunctions if no comma exists
            parts = re.split(r'(?<=[a-z])\s+(and|but|or)(?=\s+)', sentence)
            if len(parts) > 2:
                # Reconstruct with better breathing points
                # parts = [clause1, conjunction, clause2, conjunction, clause3, ...]
                phrase = parts[0]
                for i in range(1, len(parts) - 1, 2):
                    conj = parts[i]          # "and", "but", "or"
                    clause = parts[i + 1] if i + 1 < len(parts) else ""
                    # Add comma before conjunction if this is a long clause
                    phrase_words = phrase.split()
                    if len(phrase_words) > 10:
                        phrase = phrase.rstrip() + ", " + conj + " " + clause
                    else:
                        phrase = phrase + " " + conj + " " + clause
                sentence = phrase
        
        # Ensure proper punctuation at sentence end
        sentence = sentence.rstrip()
        if sentence and sentence[-1] not in {'.', '!', '?', '…', ',', ';', '—', ':'}:
            # Replace terminal ellipsis with extended pause marker
            if len(sentence) > 0:
                sentence = sentence + '...'
        elif sentence.endswith(('.', '!', '?')):
            # Replace terminal period with ellipsis for segment boundary breathing
            sentence = sentence[:-1] + '...'
        
        processed_sentences.append(sentence)
    
    # Rejoin sentences with spacing for TTS
    result = ' '.join(processed_sentences)
    
    # Add final breath marker
    return result

def _normalize_text_for_tts(text):
    """
    Safe text normalization for MeloTTS / OpenVoice.
    Only adds spaces where they are genuinely missing.
    Never removes existing spaces.
    """
    if not text or not isinstance(text, str):
        return text

    import re

    # 0. Convert LaTeX math notation to natural spoken text
    text = _latex_to_spoken(text)

    # 1. Replace em/en dashes with comma (better pause)
    text = text.replace('\u2014', ',').replace('\u2013', ',')

    # 2. Fix common LLM JSON output bugs (very conservative)
    #    Add space after punctuation ONLY if there is NO space already
    text = re.sub(r'([.!?;,])([A-Za-z])', r'\1 \2', text)

    # 3. Fix camelCase joins (e.g. "multiplexingAnd" → "multiplexing And")
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)

    # 4. Space out common 3+ letter ALL-CAPS acronyms (TCP → T C P)
    def _space_acronym(m):
        acr = m.group(0)
        if len(acr) >= 3 or acr in {'AI', 'ML', 'VR', 'AR', 'IP', 'OS', 'UI'}:
            return ' '.join(acr)
        return acr
    text = re.sub(r'\b[A-Z]{2,}\b', _space_acronym, text)

    # 5. Collapse any multiple spaces created above
    text = re.sub(r' +', ' ', text)

    return text.strip()


def smooth_autoregressive_audio(text):
    if not text or not isinstance(text, str):
        return text

    # Normalize first (em dashes, acronyms, word joins)
    text = _normalize_text_for_tts(text)

    # Then apply breathing pauses
    text = add_breathing_pauses(text)

    # Final breathing pause at end of segment
    text = text.rstrip()
    if text and text[-1] not in {'…', '.', '!', '?'}:
        text = text + '...'

    return text + " "

# ==========================================
# IMAGE CORRECTION & MANUAL ASSET SELECTION
# ==========================================
def extract_retrieve_segments(script_json):
    """
    Extract all RETRIEVE_ASSET segments from script for manual correction.
    Returns list of tuples: (segment_index, query, spoken_text)
    """
    if not script_json or str(script_json).strip() == "":
        return []
    
    try:
        script_data = json.loads(script_json)
        if isinstance(script_data, dict) and "segments" in script_data:
            segments = script_data["segments"]
        else:
            segments = script_data
        
        retrieve_segments = []
        for idx, segment in enumerate(segments):
            if segment.get('visual_mode') == 'RETRIEVE_ASSET' and segment.get('query'):
                retrieve_segments.append({
                    "segment_index": idx,
                    "query": segment.get('query', ''),
                    "spoken_text": segment.get('spoken_text', '')[:100] + "..." if len(segment.get('spoken_text', '')) > 100 else segment.get('spoken_text', '')
                })
        
        return retrieve_segments
    except Exception as e:
        print(f"Error extracting segments: {e}")
        return []

def create_image_correction_ui(script_json, pdf_images_state):
    """
    Create UI controls for manual image selection per segment.
    Returns formatted display for the image correction section.
    """
    retrieve_segments = extract_retrieve_segments(script_json)
    pdf_images = pdf_images_state if isinstance(pdf_images_state, list) else []
    
    if not retrieve_segments:
        return "ℹ️ No RETRIEVE_ASSET segments in current script."
    
    if not pdf_images:
        return "⚠️ Please generate PDF previews first (Step 1) to select manual images."
    
    # Create a simple formatted text indicating how many segments need image corrections
    info_text = f"✏️ Found {len(retrieve_segments)} image segments to correct.\n\n"
    for seg in retrieve_segments:
        info_text += f"Segment {seg['segment_index'] + 1}: {seg['query'][:50]}...\n"
    
    return info_text

def save_image_selections(image_selections_json, script_json):
    """
    Save user's manual image selections to a JSON mapping.
    image_selections_json: JSON string mapping segment indices to image paths
    Returns: confirmation message
    """
    if not image_selections_json or str(image_selections_json).strip() == "":
        return "⚠️ No image selections provided."
    
    try:
        selections = json.loads(image_selections_json)
        # Validate that selections reference valid segments
        script_retrieve_segments = extract_retrieve_segments(script_json)
        valid_indices = [s['segment_index'] for s in script_retrieve_segments]
        
        invalid_sels = [k for k in selections.keys() if int(k) not in valid_indices]
        if invalid_sels:
            return f"⚠️ Warning: Selections reference invalid segments: {invalid_sels}"
        
        return f"✅ Saved {len(selections)} image corrections. These will override auto-retrieved assets."
    except json.JSONDecodeError as e:
        return f"❌ Invalid JSON format: {e}"
    except Exception as e:
        return f"❌ Error: {e}"

def merge_manual_and_auto_assets(script_json, auto_retrieved_images, manual_image_selections_json):
    """
    Merge manual image selections with auto-retrieved assets.
    Manual selections override auto-retrieved ones for respective segments.
    Returns: merged list of (image_path, description) tuples aligned with RETRIEVE_ASSET segments
    """
    retrieve_segments = extract_retrieve_segments(script_json)
    
    # Parse manual selections
    manual_selections = {}
    if manual_image_selections_json and str(manual_image_selections_json).strip() != "":
        try:
            manual_selections = json.loads(manual_image_selections_json)
        except (json.JSONDecodeError, ValueError):
            pass
    
    # Convert manual_selections keys to integers for easier comparison
    manual_selections = {int(k): v for k, v in manual_selections.items()}
    
    merged_assets = []
    auto_idx = 0
    
    for seg in retrieve_segments:
        seg_idx = seg['segment_index']
        query = seg['query']
        
        if seg_idx in manual_selections:
            # Use manual selection
            img_path = manual_selections[seg_idx]
            merged_assets.append((img_path, f"[Manual] Query: '{query[:40]}'"))
        elif auto_idx < len(auto_retrieved_images):
            # Use auto-retrieved image
            merged_assets.append(auto_retrieved_images[auto_idx])
            auto_idx += 1
        else:
            # Fallback: use a placeholder
            merged_assets.append((PLACEHOLDER_IMAGE, f"[Missing] Query: '{query}'"))
    
    return merged_assets

def _convert_pptx_to_pdf(pptx_abs_path, pdf_output_path):
    """Convert PPTX to PDF using PowerPoint COM or LibreOffice. Returns True on success."""
    # Method 1: PowerPoint COM via comtypes
    try:
        import comtypes.client
        try:
            import pythoncom; pythoncom.CoInitialize()
        except Exception:
            pass
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            ppt.Visible = 0
            deck = ppt.Presentations.Open(pptx_abs_path, True, False, False)
            deck.SaveAs(pdf_output_path, 32)  # 32 = ppSaveAsPDF
            deck.Close()
        finally:
            try:
                ppt.Quit()
            except Exception:
                pass
        if os.path.exists(pdf_output_path):
            print("✅ PPTX→PDF via PowerPoint COM (comtypes)")
            return True
    except Exception as e:
        print(f"ℹ️ PowerPoint COM (comtypes) not available: {e}")

    # Method 2: PowerPoint COM via win32com
    try:
        import win32com.client
        try:
            import pythoncom; pythoncom.CoInitialize()
        except Exception:
            pass
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        try:
            deck = ppt.Presentations.Open(pptx_abs_path, ReadOnly=True, WithWindow=False)
            deck.SaveAs(pdf_output_path, 32)
            deck.Close()
        finally:
            try:
                ppt.Quit()
            except Exception:
                pass
        if os.path.exists(pdf_output_path):
            print("✅ PPTX→PDF via PowerPoint COM (win32com)")
            return True
    except Exception as e:
        print(f"ℹ️ PowerPoint COM (win32com) not available: {e}")

    # Method 3: LibreOffice headless
    for soffice_path in ["soffice",
                         r"C:\Program Files\LibreOffice\program\soffice.exe",
                         r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
        try:
            out_dir = os.path.dirname(pdf_output_path)
            subprocess.run(
                [soffice_path, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_abs_path],
                capture_output=True, text=True, timeout=120
            )
            lo_pdf = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_abs_path))[0] + ".pdf")
            if os.path.exists(lo_pdf):
                if lo_pdf != pdf_output_path:
                    os.replace(lo_pdf, pdf_output_path)
                print("✅ PPTX→PDF via LibreOffice")
                return True
        except Exception:
            continue

    print("⚠️ No PPTX→PDF converter available (install Microsoft Office or LibreOffice)")
    return False


def _render_pptx_to_slide_images(pptx_path, output_folder, canvas_w=1280, canvas_h=720):
    """Render each PPTX slide to a PNG image.
    Strategy: PPTX → PDF (via PowerPoint COM or LibreOffice) → images (via PyMuPDF).
    Falls back to a basic PIL renderer if no converter is available.
    Returns a dict mapping 0-based slide index to image path."""
    slide_images = {}
    abs_pptx = os.path.abspath(pptx_path)
    pdf_path = os.path.join(output_folder, "_deck_render.pdf")

    # ---- Primary: Convert PPTX → PDF, then render with PyMuPDF ----
    got_pdf = _convert_pptx_to_pdf(abs_pptx, pdf_path)
    if got_pdf:
        try:
            from PIL import Image as PILImage
            doc = fitz.open(pdf_path)
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                # Render at 2x for quality, then resize to exact canvas
                zoom = max(canvas_w / page.rect.width, canvas_h / page.rect.height) * 2
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img = img.resize((canvas_w, canvas_h), PILImage.LANCZOS)
                img_path = os.path.join(output_folder, f"pptx_slide_{page_idx:03d}.png")
                img.save(img_path)
                slide_images[page_idx] = img_path
            doc.close()
            try:
                os.remove(pdf_path)
            except OSError:
                pass
            print(f"✅ Rendered {len(slide_images)} slides via PDF→PyMuPDF")
            return slide_images
        except Exception as e:
            print(f"⚠️ PDF→image rendering failed: {e}")
            slide_images = {}

    # ---- Fallback: PIL-based renderer ----
    print("ℹ️ Falling back to PIL-based PPTX renderer")
    try:
        from pptx import Presentation as PptxPres
        from pptx.dml.color import RGBColor
        from PIL import Image as PILImage, ImageDraw, ImageFont
        import io

        prs = PptxPres(pptx_path)
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height
        sx = canvas_w / slide_w_emu
        sy = canvas_h / slide_h_emu
        pt_to_px = sy * 12700

        def _rgb_to_tuple(rgb_color, default=(51, 51, 51)):
            if rgb_color is None:
                return default
            try:
                hex_str = str(rgb_color)
                return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
            except Exception:
                return default

        _font_cache = {}
        def _get_font(size, bold=False):
            key = (size, bold)
            if key in _font_cache:
                return _font_cache[key]
            font = None
            font_names = (["arialbd.ttf", "calibrib.ttf"] if bold else []) + ["arial.ttf", "calibri.ttf", "segoeui.ttf", "verdana.ttf"]
            for fname in font_names:
                for prefix in ["", "C:/Windows/Fonts/"]:
                    try:
                        font = ImageFont.truetype(prefix + fname, size)
                        break
                    except (IOError, OSError):
                        continue
                if font:
                    break
            if font is None:
                font = ImageFont.load_default(size=size) if hasattr(ImageFont, 'load_default') else ImageFont.load_default()
            _font_cache[key] = font
            return font

        for slide_idx, slide in enumerate(prs.slides):
            bg_color = (248, 249, 250)
            try:
                bg = slide.background
                if bg.fill and bg.fill.type is not None:
                    bg_color = _rgb_to_tuple(bg.fill.fore_color.rgb, (248, 249, 250))
            except Exception:
                pass

            canvas = PILImage.new("RGB", (canvas_w, canvas_h), bg_color)
            draw = ImageDraw.Draw(canvas)

            for shape in slide.shapes:
                left = int(shape.left * sx)
                top = int(shape.top * sy)
                width = int(shape.width * sx)
                height = int(shape.height * sy)

                try:
                    if hasattr(shape, "fill"):
                        fill = shape.fill
                        if fill.type is not None and fill.fore_color and fill.fore_color.rgb:
                            color = _rgb_to_tuple(fill.fore_color.rgb, None)
                            if color:
                                draw.rectangle([left, top, left + width, top + height], fill=color)
                except Exception:
                    pass

                if shape.shape_type is not None and shape.shape_type == 13:
                    try:
                        img_blob = shape.image.blob
                        img = PILImage.open(io.BytesIO(img_blob)).convert("RGB")
                        img = img.resize((max(1, width), max(1, height)), PILImage.LANCZOS)
                        canvas.paste(img, (left, top))
                    except Exception:
                        pass
                elif hasattr(shape, "image"):
                    try:
                        img_blob = shape.image.blob
                        img = PILImage.open(io.BytesIO(img_blob)).convert("RGB")
                        img = img.resize((max(1, width), max(1, height)), PILImage.LANCZOS)
                        canvas.paste(img, (left, top))
                    except Exception:
                        pass

                if shape.has_text_frame:
                    try:
                        tf = shape.text_frame
                        y_cursor = top + 4
                        for para in tf.paragraphs:
                            text_line = para.text.strip()
                            if not text_line:
                                y_cursor += 10
                                continue
                            font_size = max(10, int(18 * pt_to_px))
                            is_bold = False
                            if para.runs:
                                run = para.runs[0]
                                if run.font.size:
                                    font_size = max(10, int(run.font.size.pt * pt_to_px))
                                is_bold = bool(run.font.bold)
                            font = _get_font(font_size, is_bold)

                            text_color = (51, 51, 51)
                            if para.runs and para.runs[0].font.color and para.runs[0].font.color.rgb:
                                text_color = _rgb_to_tuple(para.runs[0].font.color.rgb, (51, 51, 51))

                            max_text_w = max(width - 8, 50)
                            words = text_line.split()
                            lines = []
                            current_line = ""
                            for word in words:
                                test = f"{current_line} {word}".strip()
                                bbox = font.getbbox(test)
                                tw = bbox[2] - bbox[0] if bbox else len(test) * font_size * 0.6
                                if tw > max_text_w and current_line:
                                    lines.append(current_line)
                                    current_line = word
                                else:
                                    current_line = test
                            if current_line:
                                lines.append(current_line)

                            for line in lines:
                                if y_cursor < top + height:
                                    draw.text((left + 4, y_cursor), line, fill=text_color, font=font)
                                y_cursor += font_size + 4
                            y_cursor += 4
                    except Exception:
                        pass

            img_path = os.path.join(output_folder, f"pptx_slide_{slide_idx:03d}.png")
            canvas.save(img_path)
            slide_images[slide_idx] = img_path

    except Exception as e:
        print(f"⚠️ PPTX slide rendering failed: {e}")
    return slide_images
        
# ==========================================
# STEP 3: AUDIO & VIDEO SYNTHESIS (MELOTTS)
# ==========================================
def step3_render_video(script_json, ref_audio, approved_images, avatar_image, manual_selections_state=None, tts_speed=0.85, pip_enabled=True, canvas_resolution="720p (1280x720)", pexels_api_key="", presentation_mode="Classic (Avatar + Assets + B-Roll)", enable_subtitles=True, tts_engine="MeloTTS + OpenVoice (Default)", enable_breathing=True, avatar_model_root="", avatar_model_python="", melotts_ckpt_root="checkpoints_v2", avatar_backend="", progress=gr.Progress()):    
    _cancel_render.clear()  # Reset cancel flag at the start of each render

    def _render_cancelled():
        return _cancel_render.is_set()

    def _render_cleanup(*models):
        """Free GPU memory for all passed model objects."""
        for m in models:
            try:
                del m
            except Exception:
                pass
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()

    # Determine if we should use PPTX slide backgrounds
    use_pptx_slides = (presentation_mode == "Presentation Mode (PPTX Slides + Avatar PiP)")

    # Set Pexels API key at runtime if provided via UI
    global PEXELS_API_KEY
    if pexels_api_key and pexels_api_key.strip():
        PEXELS_API_KEY = pexels_api_key.strip()
    
    if not script_json or str(script_json).strip() == "":
        yield "❌ Error: Please generate a script in Tab 1 first!", None
        return

    if not ref_audio:
        yield "❌ Error: Please upload a reference audio file for Zero-Shot voice cloning!", None
        return

    yield "📝 Parsing script and initializing OpenVoice v2...", None

    _render_start_time = time.time()
    if torch.cuda.is_available():
        torch.cuda.reset_peak_memory_stats()

    try:
        script_data = json.loads(script_json)
        segments = script_data.get("segments", script_data) if isinstance(script_data, dict) else script_data
        
        if not segments:
            yield "❌ Error: Script JSON is valid, but no 'segments' were found!", None
            return
    except Exception as e:
        yield f"❌ JSON Error: {e}", None
        return

    output_folder = os.path.abspath("final_render_assets")
    os.makedirs(output_folder, exist_ok=True)

    # Clean up stale PPTX slide PNGs from previous runs so SSIM eval
    # doesn't accidentally compare old slides against the current PDF.
    for _stale in os.listdir(output_folder):
        if _stale.startswith("pptx_slide_") and _stale.endswith(".png"):
            try:
                os.remove(os.path.join(output_folder, _stale))
            except OSError:
                pass

    video_timeline = []

    # Parse canvas resolution (decoupled from avatar's native 512×512)
    canvas_w, canvas_h = CANVAS_PRESETS.get(canvas_resolution, (1280, 720))
    print(f"🖥️ Canvas: {canvas_w}×{canvas_h} | PiP: {'ON' if pip_enabled else 'OFF'}")

    device = "cuda" if torch.cuda.is_available() else "cpu"

    # Determine TTS engine via the registry (drop-down label → backend instance)
    tts_backend = get_tts_backend(tts_engine or DEFAULT_TTS_BACKEND)
    use_f5tts = (tts_backend.key == "f5-tts")  # legacy alias kept for downstream branches

    # --- 1. INITIALIZE TTS ENGINE ---
    reference_wav = ref_audio
    yield f"🗣️ Loading TTS backend '{tts_backend.name}' on {device}...", None
    try:
        tts_backend.initialize(
            device=device,
            ref_audio=reference_wav,
            melotts_ckpt_root=melotts_ckpt_root,
        )
        yield f"✅ {tts_backend.name} ready. {tts_backend.description}", None
    except FileNotFoundError as e:
        yield f"❌ TTS Initialization Error: {e}", None
        return
    except Exception as e:
        yield f"❌ TTS Initialization Error ({tts_backend.name}): {str(e)}", None
        return

    # --- 2. VISUAL ENGINE ---
    yield "👁️ Connecting to JinaCLIP Visual Database...", None
    try:
        visual_model = AutoModel.from_pretrained('jinaai/jina-clip-v1', 
                                                 trust_remote_code=True, 
                                                 low_cpu_mem_usage=True, 
                                                 _fast_init=False).to(device)
        client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
        collection = client.get_collection(name="course_assets")
    except Exception as e:
        yield f"❌ Visual Engine Error: {e}", None
        return

    # --- 3. AUTO-GENERATE PPTX SLIDE BACKGROUNDS (Presentation Mode only) ---
    pptx_slide_images = {}  # Maps segment index → slide image path
    if use_pptx_slides:
        try:
            # First check if a PPTX was already generated in Step 2 *during this session*
            # If the file is older than the current script, it's stale — remove it
            existing_pptx = os.path.join(os.path.abspath("final_render_assets"), "CUSTOM_DECK.pptx")
            if os.path.exists(existing_pptx):
                # Verify the deck is fresh: compare its mtime against render_assets_temp
                # which gets created at the start of step3. If PPTX is older, it's from
                # a previous run and must be regenerated.
                pptx_mtime = os.path.getmtime(existing_pptx)
                render_start_time = os.path.getmtime(output_folder) if os.path.exists(output_folder) else 0
                if pptx_mtime >= render_start_time:
                    yield "📊 Using existing PowerPoint deck from Step 2...", None
                    pptx_path = existing_pptx
                else:
                    yield "📊 Stale PPTX detected from a previous run — regenerating...", None
                    os.remove(existing_pptx)
                    pptx_gen = generate_pptx_deck(
                        script_json, approved_images, None,
                        ui_n_ctx=4096, ui_n_batch=512, ui_temp=0.3,
                        manual_overrides=manual_selections_state
                    )
                    pptx_status_msg, pptx_path = None, None
                    for pptx_status_msg, pptx_path in pptx_gen:
                        pass
            else:
                yield "📊 No Step 2 deck found — auto-generating PowerPoint deck for slide backgrounds...", None
                pptx_gen = generate_pptx_deck(
                    script_json, approved_images, None,
                    ui_n_ctx=4096, ui_n_batch=512, ui_temp=0.3,
                    manual_overrides=manual_selections_state
                )
                pptx_status_msg, pptx_path = None, None
                for pptx_status_msg, pptx_path in pptx_gen:
                    pass  # consume generator, keep last yielded values
            if pptx_path and os.path.exists(pptx_path):
                yield "🖼️ Rendering PPTX slides to images for video backgrounds...", None
                raw_slide_images = _render_pptx_to_slide_images(pptx_path, output_folder, canvas_w, canvas_h)
                # Map segment index → PPTX slide image
                # If slide count matches segments: 1:1 mapping
                # If fewer slides than segments: cycle through available slides
                n_slides = len(raw_slide_images)
                if n_slides > 0:
                    sorted_slides = [raw_slide_images[k] for k in sorted(raw_slide_images.keys())]
                    for seg_i in range(len(segments)):
                        pptx_slide_images[seg_i] = sorted_slides[seg_i % n_slides]
                yield f"✅ Generated {n_slides} slide backgrounds from PPTX deck → mapped to {len(segments)} segments.", None
        except Exception as e:
            yield f"⚠️ PPTX slide generation skipped: {e}", None
    else:
        yield "ℹ️ Classic mode: using avatar + retrieved assets + B-roll (no PPTX slides).", None

    # --- 4. COMPILER LOOP ---
    total_segments = len(segments)
    for i, segment in enumerate(segments):
        # --- Cancel check at the top of each segment ---
        if _render_cancelled():
            yield "⏹️ Render cancelled. Cleaning up GPU resources...", None
            tts_backend.cleanup()
            _render_cleanup(visual_model)
            yield "⏹️ Render cancelled by user. Models unloaded, VRAM freed.", None
            return

        idx = i + 1
        current_progress = i / total_segments
        text = segment.get("spoken_text", "")
        mode = segment.get("visual_mode", "GENERATE_AVATAR")
        query = segment.get("query", "")

        if not text:
            continue
        
        # Apply semantic engine text smoothing for natural audio prosody
        if enable_breathing:
            text = smooth_autoregressive_audio(text)

        yield f"🎬 Processing Segment {idx}/{len(segments)}: Generating Audio...", None

        # File paths
        temp_base = os.path.join(output_folder, f"temp_base_{idx:03d}.wav")
        final_audio = os.path.join(output_folder, f"segment_{idx:03d}.wav")
        
        try:
            # Backend-agnostic dispatch through the TTS registry.
            # MeloTTS+OpenVoice runs a two-step base→convert internally;
            # F5-TTS runs a single zero-shot clone call. The pipeline doesn't care.
            if not use_f5tts:
                yield f"🎨 Segment {idx}: Synthesising + voice-cloning...", None
            tts_backend.synthesize(
                text=text,
                output_path=final_audio,
                speed=tts_speed,
                temp_dir=output_folder,
            )

            if os.path.exists(temp_base):
                os.remove(temp_base)
                
        except Exception as e:
            yield f"❌ FATAL ERROR: Audio generation failed for segment {idx}.\nDetails: {str(e)}", None
            return

        # Visual retrieval
        # === VISUAL HANDLING (Avatar + Retrieved Asset + B-Roll PiP) ===
        visual_filename = "AVATAR_PLACEHOLDER"

        # --- Helper: progress callback factory ---
        base_progress = current_progress + (0.1 / total_segments)
        render_chunk_size = 0.8 / total_segments

        def live_avatar_progress(fraction, task_desc):
            total_fraction = base_progress + (fraction * render_chunk_size)
            progress(total_fraction, desc=f"🎥 Segment {idx}: {task_desc}")

        # ══════════════════════════════════════════════
        # MODE: GENERATE_AVATAR (full-screen talking head)
        # ══════════════════════════════════════════════
        if mode == "GENERATE_AVATAR":
            yield f"🎥 Segment {idx}: Generating talking-head avatar...", None
            raw_avatar_path = generate_talking_head_avatar(
                audio_path=final_audio,
                source_image_path=avatar_image,
                output_dir=output_folder,
                progress_callback=live_avatar_progress,
                avatar_root=avatar_model_root,
                avatar_env_python=avatar_model_python,
                avatar_backend=avatar_backend,
            )

            if raw_avatar_path:
                # Check if a PPTX slide background is available for this segment
                avatar_slide_bg = pptx_slide_images.get(i)
                if pip_enabled and avatar_slide_bg and os.path.exists(avatar_slide_bg):
                    # PiP over PPTX slide background
                    yield f"📊 Segment {idx}: Compositing avatar PiP over PPTX slide...", None
                    pip_output = os.path.join(output_folder, f"pip_{idx:03d}.mp4")
                    pip_result = create_pip_composite(
                        bg_path=avatar_slide_bg,
                        avatar_video=raw_avatar_path,
                        audio_path=final_audio,
                        output_path=pip_output,
                        canvas_w=canvas_w, canvas_h=canvas_h,
                        bg_is_video=False
                    )
                    if pip_result:
                        video_timeline.append({
                            "segment": idx,
                            "audio_file": f"segment_{idx:03d}.wav",
                            "video_file": os.path.basename(pip_result),
                            "image_file": None,
                            "script": text
                        })
                        continue

                if pip_enabled:
                    # Upscale avatar to match canvas resolution (Lanczos)
                    yield f"🔍 Segment {idx}: Upscaling avatar to {canvas_w}×{canvas_h}...", None
                    upscaled_path = os.path.join(output_folder, f"upscaled_{idx:03d}.mp4")
                    result = upscale_avatar_video(raw_avatar_path, upscaled_path, canvas_w, canvas_h)
                    visual_filename = os.path.basename(result if result else raw_avatar_path)
                else:
                    visual_filename = os.path.basename(raw_avatar_path)

                video_timeline.append({
                    "segment": idx,
                    "audio_file": f"segment_{idx:03d}.wav",
                    "video_file": visual_filename,
                    "image_file": None,
                    "script": text
                })
            else:
                # Fallback if SoulX fails — create placeholder VIDEO (not static image)
                yield f"⚠️ Segment {idx}: Avatar generation failed, using placeholder video.", None
                fallback_clip_name = f"avatar_fallback_{idx:03d}.mp4"
                fallback_clip_path = os.path.join(output_folder, fallback_clip_name)
                create_static_video_clip(PLACEHOLDER_IMAGE, final_audio, fallback_clip_path, canvas_w=canvas_w, canvas_h=canvas_h)
                video_timeline.append({
                    "segment": idx,
                    "audio_file": f"segment_{idx:03d}.wav",
                    "video_file": fallback_clip_name,
                    "image_file": None,
                    "script": text
                })
            continue

        # ══════════════════════════════════════════════
        # MODE: RETRIEVE_ASSET (diagram/chart from PDF)
        # ══════════════════════════════════════════════
        elif mode == "RETRIEVE_ASSET" and query:
            yield f"🔍 Segment {idx}: Processing '{query}'...", None
            find_image = None

            # Step 1: Check manual overrides
            if manual_selections_state and isinstance(manual_selections_state, dict):
                if i in manual_selections_state:
                    manual_path = manual_selections_state[i]
                    if os.path.exists(manual_path):
                        find_image = manual_path
                        yield f"✅ Segment {idx}: Using USER-OVERRIDE image for '{query}'...", None

            # Step 2: Auto fact-grounded retrieval
            if not find_image:
                # Try direct slide/image number lookup first (instant, no CLIP)
                direct_hit = _direct_slide_lookup(query)
                if direct_hit:
                    find_image = direct_hit
                    yield f"🔍 Segment {idx}: Direct slide lookup for '{query}'...", None

            if not find_image:
                try:
                    text_emb = visual_model.encode_text([query])
                    results = collection.query(query_embeddings=text_emb.tolist(), n_results=1)
                    if results['metadatas'] and results['metadatas'][0]:
                        find_image = results['metadatas'][0][0]['image_path']
                        yield f"🔍 Segment {idx}: Using AUTO-RETRIEVED image for '{query}'...", None
                except Exception as e:
                    print(f"Error retrieving auto asset: {e}")
                    yield f"⚠️ Segment {idx}: Failed to retrieve auto image - {str(e)}", None

            # Step 3: Use PPTX slide as background (preferred for presentation mode)
            slide_bg = pptx_slide_images.get(i)
            if slide_bg and os.path.exists(slide_bg):
                bg_image = slide_bg
                yield f"📊 Segment {idx}: Using PPTX slide as background for PiP.", None
            else:
                bg_image = find_image if find_image and os.path.exists(find_image) else PLACEHOLDER_IMAGE

            if pip_enabled:
                # --- PiP MODE: Generate avatar + compose over high-res asset ---
                yield f"🎥 Segment {idx}: Generating avatar for PiP overlay...", None
                avatar_for_pip = generate_talking_head_avatar(
                    audio_path=final_audio,
                    source_image_path=avatar_image,
                    output_dir=output_folder,
                    progress_callback=live_avatar_progress,
                    avatar_root=avatar_model_root,
                    avatar_env_python=avatar_model_python,
                    avatar_backend=avatar_backend,
                )

                if avatar_for_pip:
                    yield f"🎬 Segment {idx}: Compositing PiP ({canvas_w}×{canvas_h})...", None
                    pip_output = os.path.join(output_folder, f"pip_{idx:03d}.mp4")
                    pip_result = create_pip_composite(
                        bg_path=bg_image,
                        avatar_video=avatar_for_pip,
                        audio_path=final_audio,
                        output_path=pip_output,
                        canvas_w=canvas_w, canvas_h=canvas_h,
                        bg_is_video=False
                    )
                    if pip_result:
                        video_timeline.append({
                            "segment": idx,
                            "audio_file": f"segment_{idx:03d}.wav",
                            "video_file": os.path.basename(pip_result),
                            "image_file": None,
                            "script": text
                        })
                        continue

                # Fallback: PiP failed → static image
                yield f"⚠️ Segment {idx}: PiP compositing failed, falling back to static image.", None

            # --- STATIC MODE (non-PiP or PiP fallback) ---
            import shutil
            if bg_image and os.path.exists(bg_image) and bg_image != PLACEHOLDER_IMAGE:
                visual_filename = f"segment_{idx:03d}_visual.png"
                source_path = os.path.abspath(bg_image)
                dest_path = os.path.join(output_folder, visual_filename)
                if source_path != dest_path:
                    try:
                        shutil.copy(source_path, dest_path)
                    except Exception as e:
                        print(f"Error copying image: {e}")
                        visual_filename = os.path.basename(PLACEHOLDER_IMAGE)
            else:
                visual_filename = f"segment_{idx:03d}_placeholder.png"
                shutil.copy(PLACEHOLDER_IMAGE, os.path.join(output_folder, visual_filename))

            video_timeline.append({
                "segment": idx,
                "audio_file": f"segment_{idx:03d}.wav",
                "video_file": None,
                "image_file": visual_filename,
                "script": text
            })

        # ══════════════════════════════════════════════
        # MODE: BROLL (stock footage + avatar PiP corner)
        # ══════════════════════════════════════════════
        elif mode == "BROLL" and query:
            broll_query = segment.get("broll_query", "")
            search_display = broll_query if broll_query else query
            yield f"🎬 Segment {idx}: Searching B-roll for '{search_display}'...", None
            broll_path = find_broll_asset(
                query=query,
                broll_query=broll_query,
                narration_text=text,
                visual_model=visual_model
            )

            # Always generate animated avatar (needed for PiP or as standalone fallback)
            yield f"🎥 Segment {idx}: Generating animated avatar...", None
            avatar_video_seg = generate_talking_head_avatar(
                audio_path=final_audio,
                source_image_path=avatar_image,
                output_dir=output_folder,
                progress_callback=live_avatar_progress,
                avatar_root=avatar_model_root,
                avatar_env_python=avatar_model_python,
                avatar_backend=avatar_backend,
            )

            segment_handled = False

            # --- Priority 0 (Presentation Mode): PPTX slide + avatar PiP ---
            broll_slide_bg = pptx_slide_images.get(i)
            if not segment_handled and pip_enabled and broll_slide_bg and os.path.exists(broll_slide_bg) and avatar_video_seg:
                yield f"📊 Segment {idx}: Compositing avatar PiP over PPTX slide (presentation mode)...", None
                pip_output = os.path.join(output_folder, f"broll_pip_{idx:03d}.mp4")
                pip_result = create_pip_composite(
                    bg_path=broll_slide_bg,
                    avatar_video=avatar_video_seg,
                    audio_path=final_audio,
                    output_path=pip_output,
                    canvas_w=canvas_w, canvas_h=canvas_h,
                    bg_is_video=False
                )
                if pip_result:
                    video_timeline.append({
                        "segment": idx,
                        "audio_file": f"segment_{idx:03d}.wav",
                        "video_file": os.path.basename(pip_result),
                        "image_file": None,
                        "script": text
                    })
                    segment_handled = True

            # --- Priority 1: B-roll + avatar PiP composite ---
            if not segment_handled and pip_enabled and broll_path and avatar_video_seg:
                is_video_broll = os.path.splitext(broll_path)[1].lower() in {'.mp4', '.mov', '.avi', '.webm', '.mkv'}
                pip_output = os.path.join(output_folder, f"broll_pip_{idx:03d}.mp4")
                yield f"🎬 Segment {idx}: Compositing B-roll PiP ({canvas_w}x{canvas_h})...", None

                pip_result = create_pip_composite(
                    bg_path=broll_path,
                    avatar_video=avatar_video_seg,
                    audio_path=final_audio,
                    output_path=pip_output,
                    canvas_w=canvas_w, canvas_h=canvas_h,
                    bg_is_video=is_video_broll
                )
                if pip_result:
                    video_timeline.append({
                        "segment": idx,
                        "audio_file": f"segment_{idx:03d}.wav",
                        "video_file": os.path.basename(pip_result),
                        "image_file": None,
                        "script": text
                    })
                    segment_handled = True

            # --- Priority 2: Full-screen animated avatar (no B-roll match) ---
            if not segment_handled and avatar_video_seg:
                yield f"⚠️ Segment {idx}: No B-roll for '{query}', using animated avatar.", None
                upscaled_path = os.path.join(output_folder, f"upscaled_{idx:03d}.mp4")
                result = upscale_avatar_video(avatar_video_seg, upscaled_path, canvas_w, canvas_h)
                video_timeline.append({
                    "segment": idx,
                    "audio_file": f"segment_{idx:03d}.wav",
                    "video_file": os.path.basename(result if result else avatar_video_seg),
                    "image_file": None,
                    "script": text
                })
                segment_handled = True

            # --- Priority 3: Raw B-roll without avatar (avatar gen failed) ---
            if not segment_handled and broll_path and os.path.exists(broll_path):
                is_video_broll = os.path.splitext(broll_path)[1].lower() in {'.mp4', '.mov', '.avi', '.webm', '.mkv'}
                if is_video_broll:
                    broll_clip = os.path.join(output_folder, f"broll_{idx:03d}.mp4")
                    broll_scale = f"scale={canvas_w}:{canvas_h}:force_original_aspect_ratio=decrease:flags=lanczos,pad={canvas_w}:{canvas_h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1"
                    cmd = [
                        "ffmpeg", "-y",
                        "-stream_loop", "-1",   # Loop B-roll if shorter than audio
                        "-i", broll_path, "-i", final_audio,
                        "-vf", broll_scale,
                        "-map", "0:v", "-map", "1:a",
                        "-c:v", "libx264", "-preset", "fast", "-crf", "23",
                        "-c:a", "aac", "-b:a", "192k",
                        "-r", "25", "-pix_fmt", "yuv420p", "-shortest", broll_clip
                    ]
                    subprocess.run(cmd, capture_output=True)
                    video_timeline.append({
                        "segment": idx,
                        "audio_file": f"segment_{idx:03d}.wav",
                        "video_file": os.path.basename(broll_clip),
                        "image_file": None,
                        "script": text
                    })
                else:
                    # Static B-roll image → convert to video clip (not raw image)
                    broll_clip_name = f"broll_static_{idx:03d}.mp4"
                    broll_clip_path = os.path.join(output_folder, broll_clip_name)
                    create_static_video_clip(broll_path, final_audio, broll_clip_path, canvas_w=canvas_w, canvas_h=canvas_h)
                    video_timeline.append({
                        "segment": idx,
                        "audio_file": f"segment_{idx:03d}.wav",
                        "video_file": broll_clip_name,
                        "image_file": None,
                        "script": text
                    })
                segment_handled = True

            # --- Priority 4: Last resort — placeholder VIDEO (never a raw static image) ---
            if not segment_handled:
                yield f"⚠️ Segment {idx}: No B-roll or avatar available, generating placeholder video.", None
                placeholder_clip_name = f"placeholder_{idx:03d}.mp4"
                placeholder_clip_path = os.path.join(output_folder, placeholder_clip_name)
                create_static_video_clip(PLACEHOLDER_IMAGE, final_audio, placeholder_clip_path, canvas_w=canvas_w, canvas_h=canvas_h)
                video_timeline.append({
                    "segment": idx,
                    "audio_file": f"segment_{idx:03d}.wav",
                    "video_file": placeholder_clip_name,
                    "image_file": None,
                    "script": text
                })

    # --- 4. CLEANUP ---
    yield "💾 Saving timeline and cleaning VRAM...", None
    with open(os.path.join(output_folder, "final_timeline.json"), "w", encoding="utf-8") as f:
        json.dump(video_timeline, f, indent=4)
    
    # Offload heavy models before stitching to avoid OOM
    tts_backend.cleanup()
    del visual_model
    gc.collect()
    torch.cuda.empty_cache()

    yield "🎞️ Stitching all segments into final presentation...", None
    try:
        # Save render efficiency log for evaluation
        _render_elapsed = time.time() - _render_start_time
        _peak_vram = torch.cuda.max_memory_allocated() / (1024**2) if torch.cuda.is_available() else 0
        _render_log = {
            "total_time_s": round(_render_elapsed, 1),
            "peak_vram_mb": round(_peak_vram, 1),
            "num_segments": len(video_timeline),
        }
        try:
            with open(os.path.join(output_folder, "render_log.json"), "w") as _rl:
                json.dump(_render_log, _rl, indent=2)
        except Exception:
            pass

        # Pass absolute path to avoid file-not-found issues
        final_video_path = finalize_full_video(output_folder, video_timeline, canvas_w=canvas_w, canvas_h=canvas_h, enable_subtitles=enable_subtitles)
        # THIS IS THE ONLY YIELD THAT ACTUALLY RETURNS THE VIDEO FILE
        yield f"✅ Compilation Complete! Final video rendered.", final_video_path 
    except Exception as e:
        yield f"⚠️ Stitching failed: {str(e)}", None
# ==========================================
# POWERPOINT DECK GENERATOR (LLM-Powered)
# ==========================================
def generate_pptx_deck(script_json, fetched_assets=None, pdf_file=None,
                       ui_n_ctx=4096, ui_n_batch=512, ui_temp=0.3,
                       manual_overrides=None, llm_model_path="",
                       pptx_system_prompt=""):
    """
    Two-stage PowerPoint generation:
    1. Ask LLM to create a clean deck JSON from the narration script.
    2. Build a professional .pptx from that deck JSON with retrieved images.
    """
    if not script_json or not script_json.strip():
        yield "⚠️ No script found. Generate a script first.", None
        return

    try:
        script_data = json.loads(script_json)
        segments = script_data.get("segments", script_data) if isinstance(script_data, dict) else script_data
    except json.JSONDecodeError as e:
        yield f"⚠️ Invalid JSON: {e}", None
        return

    if not isinstance(segments, list) or len(segments) == 0:
        yield "⚠️ Script must contain a non-empty segments array.", None
        return

    # --- Stage 1: LLM generates clean deck JSON ---
    narration_summary = json.dumps(segments, indent=2)

    # Truncate narration FIRST to protect segment data, then append catalog
    max_narration_chars = int(int(ui_n_ctx) * 3.5 * 0.4)  # 40% for narration, leaving room for catalog + output
    if len(narration_summary) > max_narration_chars:
        narration_summary = narration_summary[:max_narration_chars] + "\n... (truncated)"

    # Append image catalog AFTER truncation so it doesn't eat narration content.
    # For deck generation, the relevance query is the narration itself (concatenated
    # spoken_text), since the deck only needs visuals that match what's actually said.
    deck_query_parts = []
    for seg in segments[:30]:  # cap to first 30 segments to keep query encoder fast
        if isinstance(seg, dict):
            t = seg.get("spoken_text") or seg.get("text") or ""
            if t:
                deck_query_parts.append(str(t))
    deck_query = " ".join(deck_query_parts)[:1500]
    ranked = _rank_catalog_by_relevance(METADATA_PATH, deck_query, IMAGE_CATALOG_TOP_K)
    image_catalog = _build_image_catalog(METADATA_PATH, selected_indices=ranked, ranked_note=True)
    if image_catalog:
        narration_summary += "\n" + image_catalog

    deck_json = None
    active_model = llm_model_path.strip() if llm_model_path and llm_model_path.strip() else LLM_MODEL_PATH
    try:
        import tempfile

        yield "⏳ Stage 1: Asking LLM to design slide deck...", None

        llm_config = {
            "model_path": active_model,
            "system_prompt": pptx_system_prompt.strip() if pptx_system_prompt and pptx_system_prompt.strip() else PROMPT_DECK,
            "user_content": f"Create a presentation deck from this narration script:\n{narration_summary}",
            "n_ctx": int(ui_n_ctx),
            "n_batch": int(ui_n_batch),
            "temperature": float(ui_temp),
            "max_tokens": -1,
            "top_p": 0.95,
            "stream": False,
            "json_mode": True,
            "stop": ["</s>", "<end_of_turn>", "<eos>"]
        }

        config_fd, config_path = tempfile.mkstemp(suffix=".json", prefix="llm_deck_cfg_")
        try:
            with os.fdopen(config_fd, "w", encoding="utf-8") as cf:
                json.dump(llm_config, cf)

            result = subprocess.Popen(
                [LLM_ENV_PYTHON, LLM_WORKER_SCRIPT, config_path],
                stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                text=True, encoding="utf-8"
            )
            _active_subprocesses["pptx_llm"] = result
            stdout_data, stderr_data = result.communicate()
            _active_subprocesses.pop("pptx_llm", None)

            if result.returncode != 0:
                raise RuntimeError(f"LLM worker exited with code {result.returncode}: {stderr_data[:500]}")

            # Parse the output — find the JSON line with "done" + "content"
            deck_text = None
            for line in stdout_data.strip().split("\n"):
                line = line.strip()
                if not line:
                    continue
                try:
                    data = json.loads(line)
                    if "error" in data:
                        raise RuntimeError(data["error"])
                    if "content" in data:
                        deck_text = data["content"]
                except json.JSONDecodeError:
                    continue

            if not deck_text:
                raise RuntimeError("No content returned from LLM worker")

            if not deck_text.strip().endswith("}"):
                deck_text = deck_text.strip() + "\n}"

            deck_data = json.loads(deck_text)
            deck_json = deck_data.get("slides", [])
        finally:
            try:
                os.unlink(config_path)
            except OSError:
                pass

    except Exception as e:
        print(f"⚠️ LLM deck generation failed ({e}), falling back to direct conversion.")
        yield "⚠️ LLM failed, using direct conversion fallback...", None

    # Fallback: if LLM failed, convert narration segments directly
    if not deck_json:
        deck_json = []
        for seg in segments:
            text = seg.get("spoken_text", "") or seg.get("script", "")
            sentences = re.split(r'(?<=[.!?])\s+', text.strip())
            deck_json.append({
                "title": (sentences[0][:60] + "...") if sentences and len(sentences[0]) > 60 else (sentences[0] if sentences else "Slide"),
                "bullets": [s.strip() for s in sentences[1:] if s.strip()][:6],
                "image_query": seg.get("query", "") if seg.get("visual_mode") == "RETRIEVE_ASSET" else ""
            })

    # --- Stage 1.5: Build image map from fetched assets + manual overrides ---
    asset_map = {}        # query_key -> image path
    asset_map_idx = {}    # integer index of RETRIEVE_ASSET segments -> image path
    if fetched_assets:
        asset_idx = 0
        retrieve_counter = 0
        for seg in segments:
            if seg.get("visual_mode") == "RETRIEVE_ASSET" and seg.get("query"):
                if asset_idx < len(fetched_assets):
                    item = fetched_assets[asset_idx]
                    img_path = item[0] if isinstance(item, (list, tuple)) else item
                    if img_path and os.path.exists(str(img_path)):
                        query_key = seg.get("query", "").lower().strip()
                        asset_map[query_key] = str(img_path)
                        asset_map_idx[retrieve_counter] = str(img_path)
                    asset_idx += 1
                retrieve_counter += 1

    # Apply manual overrides (take priority over auto-fetched)
    if manual_overrides and isinstance(manual_overrides, dict):
        retrieve_counter = 0
        for seg_i, seg in enumerate(segments):
            if seg.get("visual_mode") == "RETRIEVE_ASSET" and seg.get("query"):
                if seg_i in manual_overrides:
                    override_path = manual_overrides[seg_i]
                    if override_path and os.path.exists(str(override_path)):
                        query_key = seg.get("query", "").lower().strip()
                        asset_map[query_key] = str(override_path)
                        asset_map_idx[retrieve_counter] = str(override_path)
                retrieve_counter += 1

    # --- Stage 2: Build PowerPoint from deck JSON ---
    yield f"🔨 Stage 2: Building {len(deck_json)} slides...", None
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
    BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
    ACCENT_COLOR = RGBColor(0x00, 0x66, 0xCC)
    BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)
    BULLET_COLOR = RGBColor(0x44, 0x44, 0x44)

    n_images = 0

    for slide_idx, slide_data in enumerate(deck_json):
        title_text = slide_data.get("title", f"Slide {slide_idx + 1}")
        bullets = slide_data.get("bullets", [])
        img_query = slide_data.get("image_query", "").strip()

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        # Find matching image
        matched_img = None
        if img_query:
            q_lower = img_query.lower().strip()
            # Try exact match first, then fuzzy keyword match
            if q_lower in asset_map:
                matched_img = asset_map[q_lower]
            else:
                for key, path in asset_map.items():
                    if any(word in key for word in q_lower.split()[:3]):
                        matched_img = path
                        break
        # Fallback: map slide index to RETRIEVE_ASSET index for direct matching
        if matched_img is None and slide_idx in asset_map_idx:
            matched_img = asset_map_idx[slide_idx]

        has_image = matched_img is not None

        # --- Title ---
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.3), slide_w - Inches(1.2), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.LEFT

        # --- Accent line ---
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.2), Inches(3.0), Pt(4)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = ACCENT_COLOR
        line_shape.line.fill.background()

        if has_image:
            # Side-by-side: compact bullets on LEFT, large image on RIGHT
            # Gives the image ~60% of slide width so figures/tables are readable
            bullet_col_w = Inches(4.2)
            content_top = Inches(1.5)
            content_h = slide_h - content_top - Inches(0.3)

            body_box = slide.shapes.add_textbox(
                Inches(0.6), content_top, bullet_col_w, content_h
            )
            tf = body_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {bullet}"
                p.font.size = Pt(15)
                p.font.color.rgb = BULLET_COLOR
                p.space_after = Pt(6)
                p.alignment = PP_ALIGN.LEFT

            # Insert image — takes the right portion of the slide
            img_area_left = Inches(0.6) + bullet_col_w + Inches(0.3)
            img_available_w = slide_w - img_area_left - Inches(0.4)
            img_available_h = content_h

            try:
                with Image.open(matched_img) as img:
                    iw, ih = img.size
                ratio = iw / ih
                target_w = img_available_w
                target_h = int(target_w / ratio)
                if target_h > img_available_h:
                    target_h = img_available_h
                    target_w = int(target_h * ratio)
                # Center image vertically in its area
                img_top = content_top + (img_available_h - target_h) // 2
                img_left = img_area_left + (img_available_w - target_w) // 2
                slide.shapes.add_picture(matched_img, img_left, img_top, target_w, target_h)
                n_images += 1
            except Exception as e:
                print(f"⚠️ Could not insert image for slide {slide_idx + 1}: {e}")
        else:
            # Full-width bullets
            body_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.6), slide_w - Inches(1.2), slide_h - Inches(2.2)
            )
            tf = body_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = BULLET_COLOR
                p.space_after = Pt(12)
                p.alignment = PP_ALIGN.LEFT

    # Save
    output_dir = os.path.abspath("final_render_assets")
    os.makedirs(output_dir, exist_ok=True)
    pptx_path = os.path.join(output_dir, "CUSTOM_DECK.pptx")
    prs.save(pptx_path)

    n_slides = len(deck_json)
    yield f"✅ Deck generated: {n_slides} slides ({n_images} with images)", pptx_path


# ==========================================
# ONE-CLICK AUTO-PROCESS PIPELINE
# ==========================================
def auto_process_all(
    pdf_file, selected_indices,
    ui_n_ctx, ui_n_batch, ui_temp, custom_system_prompt, focus_areas, llm_model_path, long_doc_mode,
    ref_audio, avatar_img, manual_sels,
    tts_spd, pip_on, canvas_res, pexels_key, pres_mode, subs_on,
    tts_engine_choice, breathing_on,
    avatar_root_val, avatar_python_val, melotts_ckpt_val, avatar_backend_val,
    pptx_prompt,
    do_eval_script, do_eval_slide, do_eval_voice, do_eval_eff, do_eval_quiz, do_eval_mos,
    cross_auto_on, cross_api_key_val, cross_provider_val,
    cross_cold_val, cross_comp_val, cross_mos_val,
    progress=gr.Progress()
):
    """
    Sequentially runs the entire pipeline with a single click:
      1. Build visual database
      2. Generate narration script
      3. Fetch assets via fact-grounded retrieval
      4. Render video (+ optional PPTX in Presentation mode)
      5. Run evaluation
    Yields status updates to the UI.  Checks _cancel_autoprocess between stages.
    """
    _cancel_autoprocess.clear()
    _pipeline_start_time = time.time()

    def _cancelled():
        return _cancel_autoprocess.is_set()

    # ---------- pre-flight checks ----------
    if not pdf_file:
        yield ("❌ No PDF uploaded. Please upload a PDF in Step 1.", "", None, None,
               None, None, None, None, None, None, None)
        return
    if not ref_audio:
        yield ("❌ No reference audio uploaded. Please upload audio in Step 3.", "", None, None,
               None, None, None, None, None, None, None)
        return
    if not avatar_img:
        yield ("❌ No avatar headshot uploaded. Please upload a photo in Step 3.", "", None, None,
               None, None, None, None, None, None, None)
        return

    # Auto-select all pages if user hasn't picked any
    if not selected_indices:
        try:
            doc = fitz.open(pdf_file.name)
            selected_indices = list(range(len(doc)))
            doc.close()
        except Exception as e:
            yield ("❌ Could not read PDF to auto-select pages: " + str(e), "", None, None,
                   None, None, None, None, None, None, None)
            return

    # Helper to build a status tuple (auto_status, script_output, asset_gallery, video_output,
    #   eval_status, eval_chart, eval_json, eval_radar, eval_mos, eval_voice, eval_quiz, quiz_display, efficiency)
    EMPTY = (None, None, None, None, None, None, None, None, None)  # eval placeholders

    # ========== STAGE 1: Build Database ==========
    yield (f"🔄 [1/5] Building visual database...", gr.update(), gr.update(), gr.update(), *EMPTY)
    db_final_msg = ""
    for msg in build_database(pdf_file):
        db_final_msg = msg
        yield (f"🔄 [1/5] Database: {msg}", gr.update(), gr.update(), gr.update(), *EMPTY)
        if _cancelled():
            yield (f"⏹️ Cancelled during database build.", gr.update(), gr.update(), gr.update(), *EMPTY)
            return

    if "❌" in db_final_msg:
        yield (f"❌ Database build failed: {db_final_msg}", gr.update(), gr.update(), gr.update(), *EMPTY)
        return

    # ========== STAGE 2: Generate Script ==========
    if _cancelled():
        yield (f"⏹️ Cancelled before script generation.", gr.update(), gr.update(), gr.update(), *EMPTY)
        return
    yield (f"🔄 [2/5] Generating narration script...", gr.update(), gr.update(), gr.update(), *EMPTY)
    script_json = ""
    for status_msg, script_text in step1_generate_script(
        pdf_file, selected_indices, ui_n_ctx, ui_n_batch, ui_temp,
        custom_system_prompt, focus_areas, llm_model_path, long_doc_mode
    ):
        script_json = script_text
        yield (f"🔄 [2/5] Script: {status_msg}", script_text, gr.update(), gr.update(), *EMPTY)
        if _cancelled():
            _cleanup_subprocess("script_llm")
            yield (f"⏹️ Cancelled during script generation.", script_json, gr.update(), gr.update(), *EMPTY)
            return

    if not script_json or not script_json.strip():
        yield (f"❌ Script generation produced no output.", script_json, gr.update(), gr.update(), *EMPTY)
        return

    # ========== STAGE 3: Fetch Assets ==========
    if _cancelled():
        yield (f"⏹️ Cancelled before asset retrieval.", script_json, gr.update(), gr.update(), *EMPTY)
        return
    yield (f"🔄 [3/5] Retrieving visual assets...", script_json, gr.update(), gr.update(), *EMPTY)
    fetched_gallery = step2_fetch_assets(script_json)
    yield (f"✅ [3/5] Retrieved {len(fetched_gallery)} assets.", script_json, fetched_gallery, gr.update(), *EMPTY)

    if _cancelled():
        yield (f"⏹️ Cancelled after asset retrieval.", script_json, fetched_gallery, gr.update(), *EMPTY)
        return

    # ========== STAGE 4: Render Video ==========
    # Remove any stale PPTX from a previous run so Presentation Mode generates a fresh deck
    stale_pptx = os.path.join(os.path.abspath("final_render_assets"), "CUSTOM_DECK.pptx")
    if os.path.exists(stale_pptx):
        try:
            os.remove(stale_pptx)
        except OSError:
            pass
    yield (f"🔄 [4/5] Rendering video...", script_json, fetched_gallery, gr.update(), *EMPTY)
    video_path = None
    for render_msg, vid in step3_render_video(
        script_json, ref_audio, fetched_gallery, avatar_img,
        manual_sels, tts_spd, pip_on, canvas_res, pexels_key,
        pres_mode, subs_on, tts_engine_choice, breathing_on,
        avatar_root_val, avatar_python_val, melotts_ckpt_val, avatar_backend_val, progress
    ):
        video_path = vid
        yield (f"🔄 [4/5] Render: {render_msg}", script_json, fetched_gallery, vid, *EMPTY)
        if _cancelled():
            yield (f"⏹️ Cancelled during video render.", script_json, fetched_gallery, video_path, *EMPTY)
            return

    if not video_path:
        yield (f"❌ Video render produced no output.", script_json, fetched_gallery, None, *EMPTY)
        return

    # Record total pipeline time (stages 1-4, before evaluation)
    _pipeline_elapsed = time.time() - _pipeline_start_time
    _pipeline_elapsed_str = f"{int(_pipeline_elapsed // 60)}m {int(_pipeline_elapsed % 60)}s"
    print(f"\u23f1\ufe0f Pipeline stages 1-4 completed in {_pipeline_elapsed_str} ({_pipeline_elapsed:.1f}s)")

    # ========== STAGE 5: Evaluation ==========
    if _cancelled():
        yield (f"⏹️ Cancelled before evaluation.", script_json, fetched_gallery, video_path, *EMPTY)
        return
    yield (f"🔄 [5/5] Running automated evaluation...", script_json, fetched_gallery, video_path, *EMPTY)

    import plotly.graph_objects as go

    pdf_path = pdf_file.name if pdf_file else None
    ref_audio_path = ref_audio if ref_audio else ""
    eval_status_lines = []
    eval_results = {}

    gen = eval_engine.run_full_evaluation(
        script_json=script_json or "",
        pdf_path=pdf_path or "",
        reference_audio_path=ref_audio_path,
        llm_model_path="",
        run_script_fidelity=(do_eval_script and pdf_path is not None),
        run_slide_fidelity=(do_eval_slide and pdf_path is not None),
        run_voice_similarity=(do_eval_voice and bool(ref_audio_path)),
        run_efficiency=do_eval_eff,
        run_quiz=do_eval_quiz,
        run_mos=do_eval_mos,
    )

    for msg in gen:
        if msg.startswith("RESULTS:"):
            eval_results = json.loads(msg[8:])
        else:
            eval_status_lines.append(msg)
            yield (
                f"🔄 [5/5] Eval: {msg}", script_json, fetched_gallery, video_path,
                "\n".join(eval_status_lines), None, None, None, None, None, None, None, None
            )
        if _cancelled():
            yield (
                f"⏹️ Cancelled during evaluation.", script_json, fetched_gallery, video_path,
                "\n".join(eval_status_lines), None, None, None, None, None, None, None, None
            )
            return

    # Build eval charts
    metrics_chart = eval_engine.generate_metrics_bar_chart(eval_results)
    radar_chart = eval_engine.generate_radar_chart(eval_results.get("mos_scores", {}))
    voice_chart = eval_engine.generate_voice_similarity_chart(eval_results.get("voice_similarity", {}))
    quiz_chart = eval_engine.generate_quiz_chart(eval_results.get("quiz", {}))

    try:
        report_dir = eval_engine.save_evaluation_report(
            eval_results,
            {"metrics_summary": metrics_chart, "mos_radar": radar_chart,
             "voice_similarity": voice_chart, "quiz_distribution": quiz_chart},
            script_prompt=custom_system_prompt,
            pptx_prompt=pptx_prompt,
            generation_settings={"n_ctx": ui_n_ctx, "n_batch": ui_n_batch,
                                 "temperature": ui_temp, "llm_model": llm_model_path},
            script_json=script_json,
            pipeline_time_s=_pipeline_elapsed,
        )
    except Exception:
        report_dir = None

    eval_status_lines.append("✅ Evaluation complete!")

    # ========== STAGE 5b: Cross-Model Evaluation (optional) ==========
    if cross_auto_on and cross_api_key_val and cross_api_key_val.strip() and report_dir:
        if not _cancelled():
            eval_status_lines.append("\n🔀 Running cross-model evaluation...")
            yield (
                f"🔄 [5/5] Cross-model eval...", script_json, fetched_gallery, video_path,
                "\n".join(eval_status_lines), None, None, None, None, None, None, None, None
            )
            try:
                import cross_model_eval
                provider = "openai" if "OpenAI" in (cross_provider_val or "") else "gemini"
                cross_result = cross_model_eval.run_cross_eval(
                    api_key=cross_api_key_val.strip(),
                    report_dir=report_dir,
                    run_cold=cross_cold_val,
                    run_comprehension=cross_comp_val,
                    run_mos=cross_mos_val,
                    provider=provider,
                )
                eval_status_lines.append("✅ Cross-model evaluation complete!")
            except Exception as e:
                eval_status_lines.append(f"⚠️ Cross-model eval failed: {e}")
    elif cross_auto_on and (not cross_api_key_val or not cross_api_key_val.strip()):
        eval_status_lines.append("⚠️ Cross-model eval skipped — no API key provided in Process All → Cross-Model Eval section")

    quiz_display = eval_results.get("quiz", {})
    efficiency_display = eval_results.get("efficiency", {})

    yield (
        "✅ All 5 stages complete! Your video and evaluation report are ready.",
        script_json, fetched_gallery, video_path,
        "\n".join(eval_status_lines),
        metrics_chart, eval_results,
        radar_chart, eval_results.get("mos_scores", {}),
        voice_chart, quiz_chart,
        quiz_display, efficiency_display,
    )
